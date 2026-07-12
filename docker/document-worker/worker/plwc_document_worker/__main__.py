"""Creation and PDF smoke commands for the PLwC Document Worker MVP."""

from __future__ import annotations

import argparse
import base64
import importlib
import io
import json
import os
import posixpath
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path
from xml.etree.ElementTree import Element
from typing import Callable


WORK_ROOT = Path("/work")
PDF_MAX_STRUCTURAL_INPUT_FILE_SIZE = 1_500_000_000
PDF_MAX_MERGE_OUTPUT_PAGES = 100
PDF_MAX_EXTRACT_OUTPUT_PAGES = 50
PDF_MAX_SPLIT_OUTPUT_FILES = 250
PDF_MAX_ROTATE_ALL_PAGES = 100
PDF_MAX_ROTATE_SELECTED_PAGES = 50
PDF_TEXT_MAX_PAGES = 50
PDF_TEXT_MAX_CHARS = 500_000
PDF_TEXT_MAX_PREVIEW_CHARS = 1_000
ZIP_MAX_STRUCTURAL_INPUT_FILE_SIZE = 2_000_000_000
ZIP_MAX_ENTRIES = 5_000
ZIP_MAX_EXTRACTED_BYTES = 750_000_000
ZIP_MAX_SINGLE_FILE_BYTES = 250_000_000
ZIP_MAX_PATH_LENGTH = 240
ZIP_MAX_COMPRESSION_RATIO = 100
ZIP_MAX_NESTED_DEPTH = 40
OFFICE_MAX_STRUCTURAL_INPUT_FILE_SIZE = 1_000_000_000
OFFICE_MAX_TEXT_CHARS = 500_000
OFFICE_MAX_PREVIEW_CHARS = 1_000
OFFICE_MAX_XLSX_SHEETS = 10
OFFICE_MAX_XLSX_ROWS = 2_000
OFFICE_MAX_XLSX_CELLS = 50_000
OFFICE_MAX_PPTX_SLIDES = 250
OFFICE_MAX_ODF_TABLES = 10
OFFICE_MAX_ODF_ROWS = 2_000
OFFICE_MAX_ODF_CELLS = 50_000
OFFICE_MAX_XML_PART_BYTES = 100_000_000
DOCX_V2_MAX_JSON_INPUT_BYTES = 10_000_000
DOCX_V2_MAX_CONTENT_ELEMENTS = 20_000
DOCX_V2_MAX_PARAGRAPHS = 20_000
DOCX_V2_MAX_TABLES = 200
DOCX_V2_MAX_TABLE_CELLS = 100_000
DOCX_V2_MAX_IMAGES = 100
DOCX_V2_MAX_IMAGE_BYTES = 10_000_000
DOCX_V2_SCHEMA_VERSION = "docx_v2_1"
DOCX_V2_SUPPORTED_SCHEMA_VERSIONS = {DOCX_V2_SCHEMA_VERSION}
DOCX_V2_PAGE_SIZES_MM = {"A4": (210, 297), "A5": (148, 210)}
DOCX_V2_STYLE_NAMES = {"body", "title", "heading1", "heading2", "heading3", "quote", "blockquote"}
DOCX_V2_STYLE_FIELDS = {
    "font",
    "size_pt",
    "bold",
    "italic",
    "alignment",
    "line_spacing",
    "space_before_pt",
    "space_after_pt",
    "first_line_indent_mm",
    "page_break_before",
}
DOCX_V2_TEXT_TYPES = {"title", "heading1", "heading2", "heading3", "paragraph", "quote", "blockquote"}
DOCX_V2_CONTENT_TYPES = DOCX_V2_TEXT_TYPES | {"bullet_list", "numbered_list", "page_break", "image", "table"}
DOCX_V2_RUN_FIELDS = {"text", "bold", "italic", "underline"}
DOCX_V2_PROTECTED_SEGMENTS = {"profile", "profiles", "governance"}
DOCX_V2_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}
XLSX_V2_MAX_JSON_INPUT_BYTES = 10_000_000
XLSX_V2_MAX_SHEETS = 50
XLSX_V2_MAX_ROWS_PER_SHEET = 100_000
XLSX_V2_MAX_COLUMNS_PER_SHEET = 256
XLSX_V2_MAX_CELLS_PER_SHEET = 500_000
XLSX_V2_MAX_TOTAL_CELLS = 1_000_000
XLSX_V2_MAX_MERGED_RANGES = 1_000
XLSX_V2_SCHEMA_VERSION = "xlsx_v2_1"
XLSX_V2_SUPPORTED_SCHEMA_VERSIONS = {XLSX_V2_SCHEMA_VERSION}
XLSX_V2_TOP_LEVEL_FIELDS = {"schema_version", "workbook", "sheets"}
XLSX_V2_WORKBOOK_FIELDS = {"title", "author"}
XLSX_V2_SHEET_FIELDS = {"name", "freeze_panes", "auto_filter", "column_widths", "row_heights", "merge_cells", "rows"}
XLSX_V2_FORMAT_FIELDS = {
    "bold",
    "italic",
    "underline",
    "font_size",
    "font_color",
    "fill_color",
    "alignment",
    "vertical_alignment",
    "number_format",
    "border",
    "wrap_text",
}
XLSX_V2_CELL_FIELDS = {"value", "formula"} | XLSX_V2_FORMAT_FIELDS
XLSX_V2_ALIGNMENTS = {"left", "center", "right"}
XLSX_V2_VERTICAL_ALIGNMENTS = {"top", "middle", "bottom"}
XLSX_V2_BORDERS = {"none", "thin", "medium", "thick"}
XLSX_V2_UNSAFE_FORMULA_TOKENS = ("WEBSERVICE(", "FILTERXML(", "HYPERLINK(", "RTD(", "DDE", "CALL(", "REGISTER.ID(", "EXEC(")
PPTX_V2_MAX_JSON_INPUT_BYTES = 10_000_000
PPTX_V2_MAX_SLIDES = 250
PPTX_V2_MAX_CONTENT_ELEMENTS_PER_SLIDE = 50
PPTX_V2_MAX_BULLET_ITEMS_PER_SLIDE = 200
PPTX_V2_MAX_TABLES_PER_SLIDE = 5
PPTX_V2_MAX_TABLE_CELLS_PER_SLIDE = 2_000
PPTX_V2_MAX_IMAGES = 100
PPTX_V2_MAX_IMAGE_BYTES = 10_000_000
PPTX_V2_SCHEMA_VERSION = "pptx_v2_1"
PPTX_V2_SUPPORTED_SCHEMA_VERSIONS = {PPTX_V2_SCHEMA_VERSION}
PPTX_V2_TOP_LEVEL_FIELDS = {"schema_version", "presentation", "slides"}
PPTX_V2_PRESENTATION_FIELDS = {"title", "author", "slide_size"}
PPTX_V2_SLIDE_FIELDS = {"layout", "title", "subtitle", "body", "content", "image", "notes"}
PPTX_V2_LAYOUTS = {"title", "content", "section_header", "image", "blank"}
PPTX_V2_CONTENT_TYPES = {"bullets", "paragraph", "table", "image", "text_box"}
PPTX_V2_BULLETS_FIELDS = {"type", "items"}
PPTX_V2_BULLET_ITEM_FIELDS = {"text", "level", "bold", "italic", "font_size", "color"}
PPTX_V2_PARAGRAPH_FIELDS = {"type", "text", "bold", "italic", "font_size", "color", "alignment"}
PPTX_V2_TABLE_FIELDS = {"type", "rows", "header_row"}
PPTX_V2_IMAGE_ELEMENT_FIELDS = {"type", "path", "width_mm", "height_mm", "align"}
PPTX_V2_TEXT_BOX_FIELDS = {
    "type",
    "text",
    "left_mm",
    "top_mm",
    "width_mm",
    "height_mm",
    "bold",
    "italic",
    "font_size",
    "color",
    "alignment",
}
PPTX_V2_SLIDE_IMAGE_FIELDS = {"path", "width_mm", "height_mm", "align"}
PPTX_V2_ALIGNMENTS = {"left", "center", "right"}
PPTX_V2_BULLET_LEVELS = {0, 1, 2, 3}
PPTX_V2_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}
PPTX_V2_SLIDE_SIZE_PRESETS_MM = {
    "16:9": (338, 190),
    "4:3": (254, 190),
    "A4_portrait": (210, 297),
}
PPTX_V2_CUSTOM_SIZE_FIELDS = {"name", "width_mm", "height_mm"}
PPTX_V2_CUSTOM_SIZE_MIN_MM = 100
PPTX_V2_CUSTOM_SIZE_MAX_MM = 1200
PPTX_V2_PROTECTED_SEGMENTS = {"profile", "profiles", "governance"}
PDF_V2_MAX_JSON_INPUT_BYTES = 10_000_000
PDF_V2_MAX_CONTENT_ELEMENTS = 20_000
PDF_V2_MAX_PARAGRAPHS = 20_000
PDF_V2_MAX_TABLES = 200
PDF_V2_MAX_TABLE_CELLS = 100_000
PDF_V2_MAX_IMAGES = 100
PDF_V2_MAX_IMAGE_BYTES = 10_000_000
PDF_V2_SCHEMA_VERSION = "pdf_v2_1"
PDF_V2_SUPPORTED_SCHEMA_VERSIONS = {PDF_V2_SCHEMA_VERSION}
PDF_V2_TOP_LEVEL_FIELDS = {"schema_version", "document", "page", "styles", "content"}
PDF_V2_DOCUMENT_FIELDS = {"title", "author", "language"}
PDF_V2_PAGE_FIELDS = {"size", "orientation", "margins_mm"}
PDF_V2_MARGIN_FIELDS = {"top", "bottom", "left", "right"}
PDF_V2_STYLE_NAMES = {"body", "title", "heading1", "heading2", "heading3", "quote", "blockquote"}
PDF_V2_STYLE_FIELDS = {
    "font",
    "size_pt",
    "bold",
    "italic",
    "underline",
    "alignment",
    "line_spacing",
    "space_before_pt",
    "space_after_pt",
    "first_line_indent_mm",
    "page_break_before",
}
PDF_V2_TEXT_TYPES = {"title", "heading1", "heading2", "heading3", "paragraph", "quote", "blockquote"}
PDF_V2_CONTENT_TYPES = PDF_V2_TEXT_TYPES | {"bullet_list", "numbered_list", "page_break", "image", "table"}
PDF_V2_RUN_FIELDS = {"text", "bold", "italic", "underline"}
PDF_V2_TEXT_ELEMENT_FIELDS = {"type", "text", "runs", "style", "style_overrides"}
PDF_V2_LIST_FIELDS = {"type", "items", "style_overrides"}
PDF_V2_IMAGE_FIELDS = {"type", "path", "width_mm", "height_mm", "align"}
PDF_V2_TABLE_FIELDS = {"type", "rows", "header_row"}
PDF_V2_PAGE_SIZES_MM = {"A4": (210, 297), "A5": (148, 210)}
PDF_V2_ALIGNMENTS = {"left", "center", "right", "justify"}
PDF_V2_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}
PDF_V2_CUSTOM_SIZE_FIELDS = {"name", "width_mm", "height_mm"}
PDF_V2_CUSTOM_SIZE_MIN_MM = 100
PDF_V2_CUSTOM_SIZE_MAX_MM = 1200
PDF_V2_PROTECTED_SEGMENTS = {"profile", "profiles", "governance"}
READ_IMAGE_SUPPORTED_INPUT_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
READ_IMAGE_SUPPORTED_FORMATS = ["png", "jpg", "jpeg", "webp", "gif"]
READ_IMAGE_OUTPUT_FORMATS = {"png", "jpeg", "webp"}
READ_IMAGE_MEDIA_TYPES = {"png": "image/png", "jpeg": "image/jpeg", "webp": "image/webp"}
READ_IMAGE_DEFAULT_MAX_SIZE_BYTES = 2_048 * 1024
READ_IMAGE_HARD_MAX_SIZE_BYTES = 5_120 * 1024
EDIT_DOCX_MAX_INPUT_BYTES = 1_000_000_000
EDIT_DOCX_MAX_XML_PART_BYTES = 100_000_000
EDIT_DOCX_MAX_SCANNED_PARAGRAPHS = 50_000
EDIT_DOCX_STYLE_MAP = {
    "body": "Normal",
    "title": "Title",
    "heading1": "Heading 1",
    "heading2": "Heading 2",
    "heading3": "Heading 3",
    "quote": "Quote",
    "blockquote": "Quote",
}
IMPORT_PROBES = {
    "pypdf": "pypdf",
    "python-docx": "docx",
    "openpyxl": "openpyxl",
    "python-pptx": "pptx",
    "reportlab": "reportlab",
    "weasyprint": "weasyprint",
    "odfpy": "odf",
    "pandas": "pandas",
    "defusedxml": "defusedxml",
    "Pillow": "PIL",
}


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(prog="plwc_document_worker")
    subparsers = parser.add_subparsers(dest="command", required=True)
    subparsers.add_parser("probe")

    for command, suffix in {
        "create-docx": ".docx",
        "create-xlsx": ".xlsx",
        "create-pptx": ".pptx",
        "create-pdf": ".pdf",
    }.items():
        sub = subparsers.add_parser(command)
        sub.add_argument("--output", required=True)
        sub.add_argument("--title", default="PLwC Document Worker Smoke")
        sub.add_argument("--body", default="Document Worker MVP creation smoke test.")
        sub.add_argument("--content-json", default="{}")
        sub.add_argument("--input-json", default="")
        sub.add_argument("--max-json-input-bytes", type=int, default=DOCX_V2_MAX_JSON_INPUT_BYTES)
        sub.set_defaults(expected_suffix=suffix)

    inspect_pdf = subparsers.add_parser("inspect-pdf")
    inspect_pdf.add_argument("--input", required=True)
    inspect_pdf.add_argument("--max-input-bytes", type=int, default=PDF_MAX_STRUCTURAL_INPUT_FILE_SIZE)

    merge_pdf = subparsers.add_parser("merge-pdf")
    merge_pdf.add_argument("--inputs-json", required=True)
    merge_pdf.add_argument("--output", required=True)
    merge_pdf.add_argument("--max-input-bytes", type=int, default=PDF_MAX_STRUCTURAL_INPUT_FILE_SIZE)
    merge_pdf.add_argument("--max-pages", type=int, default=PDF_MAX_MERGE_OUTPUT_PAGES)

    split_pdf = subparsers.add_parser("split-pdf")
    split_pdf.add_argument("--input", required=True)
    split_pdf.add_argument("--output-dir", required=True)
    split_pdf.add_argument("--max-input-bytes", type=int, default=PDF_MAX_STRUCTURAL_INPUT_FILE_SIZE)
    split_pdf.add_argument("--max-output-files", type=int, default=PDF_MAX_SPLIT_OUTPUT_FILES)

    extract_pdf = subparsers.add_parser("extract-pdf")
    extract_pdf.add_argument("--input", required=True)
    extract_pdf.add_argument("--output", required=True)
    extract_pdf.add_argument("--pages-json", required=True)
    extract_pdf.add_argument("--max-input-bytes", type=int, default=PDF_MAX_STRUCTURAL_INPUT_FILE_SIZE)
    extract_pdf.add_argument("--max-pages", type=int, default=PDF_MAX_EXTRACT_OUTPUT_PAGES)

    rotate_pdf = subparsers.add_parser("rotate-pdf")
    rotate_pdf.add_argument("--input", required=True)
    rotate_pdf.add_argument("--output", required=True)
    rotate_pdf.add_argument("--rotation", type=int, required=True)
    rotate_pdf.add_argument("--pages-json", default="[]")
    rotate_pdf.add_argument("--max-input-bytes", type=int, default=PDF_MAX_STRUCTURAL_INPUT_FILE_SIZE)
    rotate_pdf.add_argument("--max-all-pages", type=int, default=PDF_MAX_ROTATE_ALL_PAGES)
    rotate_pdf.add_argument("--max-selected-pages", type=int, default=PDF_MAX_ROTATE_SELECTED_PAGES)

    extract_pdf_text = subparsers.add_parser("extract-pdf-text")
    extract_pdf_text.add_argument("--input", required=True)
    extract_pdf_text.add_argument("--output", default="")
    extract_pdf_text.add_argument("--format", choices=("text", "json"), default="text")
    extract_pdf_text.add_argument("--pages-json", default="[]")
    extract_pdf_text.add_argument("--max-input-bytes", type=int, default=PDF_MAX_STRUCTURAL_INPUT_FILE_SIZE)
    extract_pdf_text.add_argument("--max-pages", type=int, default=PDF_TEXT_MAX_PAGES)
    extract_pdf_text.add_argument("--max-chars", type=int, default=PDF_TEXT_MAX_CHARS)
    extract_pdf_text.add_argument("--include-preview", action="store_true")
    extract_pdf_text.add_argument("--max-preview-chars", type=int, default=PDF_TEXT_MAX_PREVIEW_CHARS)

    inspect_zip = subparsers.add_parser("inspect-zip")
    inspect_zip.add_argument("--input", required=True)
    inspect_zip.add_argument("--max-input-bytes", type=int, default=ZIP_MAX_STRUCTURAL_INPUT_FILE_SIZE)
    inspect_zip.add_argument("--max-entries", type=int, default=ZIP_MAX_ENTRIES)
    inspect_zip.add_argument("--max-total-uncompressed-bytes", type=int, default=ZIP_MAX_EXTRACTED_BYTES)
    inspect_zip.add_argument("--max-single-file-bytes", type=int, default=ZIP_MAX_SINGLE_FILE_BYTES)
    inspect_zip.add_argument("--max-path-length", type=int, default=ZIP_MAX_PATH_LENGTH)
    inspect_zip.add_argument("--max-compression-ratio", type=int, default=ZIP_MAX_COMPRESSION_RATIO)
    inspect_zip.add_argument("--max-depth", type=int, default=ZIP_MAX_NESTED_DEPTH)

    extract_zip = subparsers.add_parser("extract-zip")
    extract_zip.add_argument("--input", required=True)
    extract_zip.add_argument("--output-dir", required=True)
    extract_zip.add_argument("--max-input-bytes", type=int, default=ZIP_MAX_STRUCTURAL_INPUT_FILE_SIZE)
    extract_zip.add_argument("--max-entries", type=int, default=ZIP_MAX_ENTRIES)
    extract_zip.add_argument("--max-total-uncompressed-bytes", type=int, default=ZIP_MAX_EXTRACTED_BYTES)
    extract_zip.add_argument("--max-single-file-bytes", type=int, default=ZIP_MAX_SINGLE_FILE_BYTES)
    extract_zip.add_argument("--max-path-length", type=int, default=ZIP_MAX_PATH_LENGTH)
    extract_zip.add_argument("--max-compression-ratio", type=int, default=ZIP_MAX_COMPRESSION_RATIO)
    extract_zip.add_argument("--max-depth", type=int, default=ZIP_MAX_NESTED_DEPTH)

    create_zip = subparsers.add_parser("create-zip")
    create_zip.add_argument("--inputs-json", required=True)
    create_zip.add_argument("--output", required=True)
    create_zip.add_argument("--max-files", type=int, default=ZIP_MAX_ENTRIES)
    create_zip.add_argument("--max-total-input-bytes", type=int, default=ZIP_MAX_EXTRACTED_BYTES)
    create_zip.add_argument("--max-single-file-bytes", type=int, default=ZIP_MAX_SINGLE_FILE_BYTES)
    create_zip.add_argument("--max-path-length", type=int, default=ZIP_MAX_PATH_LENGTH)
    create_zip.add_argument("--max-depth", type=int, default=ZIP_MAX_NESTED_DEPTH)

    for command, suffix in {
        "inspect-docx": ".docx",
        "inspect-xlsx": ".xlsx",
        "inspect-pptx": ".pptx",
        "inspect-odt": ".odt",
        "inspect-ods": ".ods",
        "inspect-odp": ".odp",
    }.items():
        inspect_office = subparsers.add_parser(command)
        inspect_office.add_argument("--input", required=True)
        inspect_office.add_argument("--max-input-bytes", type=int, default=OFFICE_MAX_STRUCTURAL_INPUT_FILE_SIZE)
        inspect_office.add_argument("--max-xml-part-bytes", type=int, default=OFFICE_MAX_XML_PART_BYTES)
        inspect_office.set_defaults(expected_suffix=suffix)

    for command, suffix in {
        "extract-docx-text": ".docx",
        "extract-pptx-text": ".pptx",
        "extract-odt-text": ".odt",
        "extract-odp-text": ".odp",
    }.items():
        extract_office_text = subparsers.add_parser(command)
        extract_office_text.add_argument("--input", required=True)
        extract_office_text.add_argument("--output", default="")
        extract_office_text.add_argument("--format", choices=("text", "json"), default="text")
        extract_office_text.add_argument("--pages-json", default="[]")
        extract_office_text.add_argument("--max-input-bytes", type=int, default=OFFICE_MAX_STRUCTURAL_INPUT_FILE_SIZE)
        extract_office_text.add_argument("--max-chars", type=int, default=OFFICE_MAX_TEXT_CHARS)
        extract_office_text.add_argument("--include-preview", action="store_true")
        extract_office_text.add_argument("--max-preview-chars", type=int, default=OFFICE_MAX_PREVIEW_CHARS)
        extract_office_text.add_argument("--max-xml-part-bytes", type=int, default=OFFICE_MAX_XML_PART_BYTES)
        extract_office_text.add_argument("--include-notes", action="store_true")
        extract_office_text.set_defaults(expected_suffix=suffix)

    for command, suffix in {"extract-xlsx-data": ".xlsx", "extract-ods-data": ".ods"}.items():
        extract_office_data = subparsers.add_parser(command)
        extract_office_data.add_argument("--input", required=True)
        extract_office_data.add_argument("--output", default="")
        extract_office_data.add_argument("--sheets-json", default="[]")
        extract_office_data.add_argument("--ranges-json", default="{}")
        extract_office_data.add_argument("--max-input-bytes", type=int, default=OFFICE_MAX_STRUCTURAL_INPUT_FILE_SIZE)
        extract_office_data.add_argument("--max-rows", type=int, default=OFFICE_MAX_XLSX_ROWS)
        extract_office_data.add_argument("--max-cells", type=int, default=OFFICE_MAX_XLSX_CELLS)
        extract_office_data.add_argument("--include-preview", action="store_true")
        extract_office_data.add_argument("--max-preview-chars", type=int, default=OFFICE_MAX_PREVIEW_CHARS)
        extract_office_data.add_argument("--max-xml-part-bytes", type=int, default=OFFICE_MAX_XML_PART_BYTES)
        extract_office_data.set_defaults(expected_suffix=suffix)

    read_image = subparsers.add_parser("read-image")
    read_image.add_argument("--input", required=True)
    read_image.add_argument("--max-size-bytes", type=int, default=READ_IMAGE_DEFAULT_MAX_SIZE_BYTES)
    read_image.add_argument("--hard-max-size-bytes", type=int, default=READ_IMAGE_HARD_MAX_SIZE_BYTES)
    read_image.add_argument("--resize-to", default="")
    read_image.add_argument("--format", choices=sorted(READ_IMAGE_OUTPUT_FORMATS), default="png")

    edit_docx = subparsers.add_parser("edit-docx")
    edit_docx.add_argument("--input", required=True)
    edit_docx.add_argument("--output", required=True)
    edit_docx.add_argument("--edits-json", required=True)
    edit_docx.add_argument("--max-input-bytes", type=int, default=EDIT_DOCX_MAX_INPUT_BYTES)
    edit_docx.add_argument("--max-xml-part-bytes", type=int, default=EDIT_DOCX_MAX_XML_PART_BYTES)

    args = parser.parse_args(argv)
    try:
        if args.command == "probe":
            return _probe()
        if args.command == "inspect-pdf":
            return _inspect_pdf(args)
        if args.command == "merge-pdf":
            return _merge_pdf(args)
        if args.command == "split-pdf":
            return _split_pdf(args)
        if args.command == "extract-pdf":
            return _extract_pdf(args)
        if args.command == "rotate-pdf":
            return _rotate_pdf(args)
        if args.command == "extract-pdf-text":
            return _extract_pdf_text(args)
        if args.command == "inspect-zip":
            return _inspect_zip(args)
        if args.command == "extract-zip":
            return _extract_zip(args)
        if args.command == "create-zip":
            return _create_zip(args)
        if args.command in {"inspect-docx", "inspect-xlsx", "inspect-pptx", "inspect-odt", "inspect-ods", "inspect-odp"}:
            return _inspect_office(args)
        if args.command in {"extract-docx-text", "extract-pptx-text", "extract-odt-text", "extract-odp-text"}:
            return _extract_office_text(args)
        if args.command in {"extract-xlsx-data", "extract-ods-data"}:
            return _extract_office_data(args)
        if args.command == "read-image":
            return _read_image(args)
        if args.command == "edit-docx":
            return _edit_docx(args)
        output = _validate_output_path(args.output, args.expected_suffix)
        operations: dict[str, Callable[[Path, str, str, dict[str, object]], dict[str, object] | None]] = {
            "create-docx": _create_docx,
            "create-xlsx": _create_xlsx,
            "create-pptx": _create_pptx,
            "create-pdf": _create_pdf,
        }
        content, input_metadata = _creation_content_from_args(args)
        result_metadata = operations[args.command](output, args.title, args.body, content) or {}
        payload = {
            "ok": True,
            "operation": args.command,
            "output_path": str(output),
            "file_size": output.stat().st_size,
        }
        payload.update(input_metadata)
        payload.update(result_metadata)
        _write_json(payload)
        return 0
    except Exception as exc:  # pragma: no cover - exercised in the worker image.
        _write_json({"ok": False, "operation": args.command, "error": str(exc), "error_category": _error_category(exc)})
        return 1


def _error_category(exc: Exception) -> str:
    if isinstance(exc, (ImportError, ModuleNotFoundError)):
        return "engine_missing"
    if not isinstance(exc, ValueError):
        return "worker_error"
    message = str(exc).casefold()
    if "destination_file_exists" in message:
        return "destination_file_exists"
    if "already exists" in message:
        return "overwrite_denied" if "zip" in message or "directory" in message else "output_exists"
    if "encrypted zip" in message or "password-protected" in message:
        return "unsupported_encrypted_zip"
    if "symlink" in message:
        return "symlink_not_supported"
    if "absolute path" in message or "drive-qualified" in message or "unc path" in message or "null byte" in message:
        return "invalid_path"
    if "zip slip" in message or "parent traversal" in message:
        return "zip_slip_detected"
    if "must use the" in message and "extension" in message:
        return "unsupported_extension"
    if "unsupported_format" in message:
        return "validation_error"
    if "file_too_large" in message:
        return "policy_violation"
    if "resize_unavailable" in message:
        return "resize_unavailable"
    if "path_traversal_denied" in message:
        return "path_traversal_denied"
    if "macro-enabled" in message or "macro" in message:
        return "unsupported_macro_document"
    if "external relationship" in message or "external link" in message:
        return "unsafe_external_relationship"
    if "malformed" in message or "bad zip" in message or "not a zip" in message:
        return "malformed_document"
    if "must stay under /work" in message:
        return "workspace_boundary_rejected"
    if (
        "out of range" in message
        or "page 0" in message
        or "1-based" in message
        or "must be a positive integer" in message
        or "pages-json" in message
        or "must decode" in message
        or "validation_error" in message
        or "unsupported field" in message
    ):
        return "validation_error"
    if "protected path" in message:
        return "protected_path_rejected"
    if "exceeds" in message or "limit is" in message:
        return "limit_exceeded"
    return "worker_error"


def _probe() -> int:
    imported: dict[str, bool] = {}
    errors: dict[str, str] = {}
    for package_name, import_name in IMPORT_PROBES.items():
        try:
            importlib.import_module(import_name)
            imported[package_name] = True
        except Exception as exc:  # pragma: no cover - depends on worker image contents.
            imported[package_name] = False
            errors[package_name] = str(exc)
    ok = all(imported.values())
    payload = {"ok": ok, "operation": "probe", "imports": imported}
    if errors:
        payload["errors"] = errors
    _write_json(payload)
    return 0 if ok else 1


def _validate_output_path(raw_output: str, expected_suffix: str) -> Path:
    output = Path(raw_output)
    if not output.is_absolute():
        raise ValueError("Output path must be absolute inside the worker.")
    if output.suffix.lower() != expected_suffix:
        raise ValueError(f"Output path must use the {expected_suffix} extension.")
    resolved_parent = output.parent.resolve(strict=False)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved_parent != work_root and work_root not in resolved_parent.parents:
        raise ValueError("Output path must stay under /work.")
    output.parent.mkdir(parents=True, exist_ok=True)
    return output


def _validate_input_path(raw_input: str, expected_suffix: str = ".pdf") -> Path:
    input_path = Path(raw_input)
    if not input_path.is_absolute():
        raise ValueError("Input path must be absolute inside the worker.")
    if input_path.suffix.lower() != expected_suffix:
        raise ValueError(f"Input path must use the {expected_suffix} extension.")
    resolved = input_path.resolve(strict=True)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved != work_root and work_root not in resolved.parents:
        raise ValueError("Input path must stay under /work.")
    if not resolved.is_file():
        raise ValueError("Input path must be a file.")
    return resolved


def _validate_image_input_path(raw_input: str) -> Path:
    input_path = Path(raw_input)
    if not input_path.is_absolute():
        raise ValueError("Input path must be absolute inside the worker.")
    if input_path.suffix.lower() not in READ_IMAGE_SUPPORTED_INPUT_SUFFIXES:
        raise ValueError("unsupported_format")
    resolved = input_path.resolve(strict=True)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved != work_root and work_root not in resolved.parents:
        raise ValueError("Input path must stay under /work.")
    if not resolved.is_file():
        raise ValueError("Input path must be a file.")
    return resolved


def _validate_zip_source_path(raw_input: str) -> Path:
    source = Path(raw_input)
    if not source.is_absolute():
        raise ValueError("ZIP source path must be absolute inside the worker.")
    resolved = source.resolve(strict=True)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved != work_root and work_root not in resolved.parents:
        raise ValueError("ZIP source path must stay under /work.")
    if source.is_symlink() or resolved.is_symlink():
        raise ValueError("ZIP symlink sources are not supported in the MVP.")
    return resolved


def _validate_output_dir(raw_output_dir: str) -> Path:
    output_dir = Path(raw_output_dir)
    if not output_dir.is_absolute():
        raise ValueError("Output directory must be absolute inside the worker.")
    resolved_parent = output_dir.parent.resolve(strict=False)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved_parent != work_root and work_root not in resolved_parent.parents:
        raise ValueError("Output directory must stay under /work.")
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir


def _validate_zip_output_dir(raw_output_dir: str) -> Path:
    output_dir = Path(raw_output_dir)
    if not output_dir.is_absolute():
        raise ValueError("ZIP output directory must be absolute inside the worker.")
    resolved_parent = output_dir.parent.resolve(strict=False)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved_parent != work_root and work_root not in resolved_parent.parents:
        raise ValueError("ZIP output directory must stay under /work.")
    # Note: pre-existing output_dir is no longer rejected here. The extractor
    # handles three cases: non-existent (atomic rename), existing-empty/safe
    # (merge), and existing-with-collision (loud destination_file_exists).
    return output_dir


def _check_input_size(path: Path, max_bytes: int) -> int:
    size = path.stat().st_size
    if size > max_bytes:
        raise ValueError(f"Input file exceeds structural safety limit of {max_bytes} bytes.")
    return size


def _read_image(args: argparse.Namespace) -> int:
    try:
        input_path = _validate_image_input_path(args.input)
    except Exception as exc:
        suffix = Path(str(args.input)).suffix.lower()
        payload: dict[str, object] = {
            "ok": False,
            "operation": "read_image",
            "error": "unsupported_format" if str(exc) == "unsupported_format" else str(exc),
            "error_category": "validation_error",
        }
        if str(exc) == "unsupported_format":
            payload["detected_extension"] = suffix
            payload["supported_formats"] = READ_IMAGE_SUPPORTED_FORMATS
        _write_json(payload)
        return 1

    original_size = input_path.stat().st_size
    hard_max_size_bytes = min(int(args.hard_max_size_bytes), READ_IMAGE_HARD_MAX_SIZE_BYTES)
    max_size_bytes = min(int(args.max_size_bytes), hard_max_size_bytes)
    suggestion = "Use resize_to parameter to reduce image size before transfer, e.g. resize_to='1024x1024' or resize_to='50%'"
    if original_size > hard_max_size_bytes:
        _write_json(
            {
                "ok": False,
                "operation": "read_image",
                "error": "file_too_large",
                "error_category": "policy_violation",
                "file_size_bytes": original_size,
                "max_size_bytes": hard_max_size_bytes,
                "suggestion": suggestion,
            }
        )
        return 1

    try:
        from PIL import Image, ImageOps
    except Exception:
        _write_json(
            {
                "ok": False,
                "operation": "read_image",
                "error": "resize_unavailable",
                "error_category": "resize_unavailable",
                "file_size_bytes": original_size,
                "max_size_bytes": max_size_bytes,
                "fallback_available": original_size <= max_size_bytes,
                "suggestion": "Pillow is required for image resize/conversion. Retry without resize_to only if the original file is within the size limit.",
            }
        )
        return 1

    with Image.open(input_path) as opened:
        if getattr(opened, "is_animated", False):
            opened.seek(0)
        image = ImageOps.exif_transpose(opened)
        output_format = str(args.format).casefold()
        resized = False
        resize_to = str(args.resize_to or "").strip().casefold()
        if resize_to:
            image = _resize_image(image, resize_to)
            resized = True
        encoded = _encode_image(image, output_format)

    if len(encoded) > max_size_bytes:
        _write_json(
            {
                "ok": False,
                "operation": "read_image",
                "error": "file_too_large",
                "error_category": "policy_violation",
                "file_size_bytes": len(encoded),
                "original_file_size_bytes": original_size,
                "max_size_bytes": max_size_bytes,
                "suggestion": suggestion,
            }
        )
        return 1

    _write_json(
        {
            "ok": True,
            "operation": "read_image",
            "input_path": str(input_path),
            "media_type": READ_IMAGE_MEDIA_TYPES[output_format],
            "width_px": image.width,
            "height_px": image.height,
            "file_size_bytes": len(encoded),
            "original_file_size_bytes": original_size,
            "resized": resized,
            "format": output_format,
            "image_data": base64.b64encode(encoded).decode("ascii"),
        }
    )
    return 0


def _resize_image(image: object, resize_to: str):
    width, height = image.size
    if resize_to.endswith("%"):
        ratio = int(resize_to[:-1]) / 100
        new_size = (max(1, int(width * ratio)), max(1, int(height * ratio)))
        return image.resize(new_size)
    width_text, height_text = resize_to.split("x", 1)
    box = (int(width_text), int(height_text))
    resized = image.copy()
    resized.thumbnail(box)
    return resized


def _encode_image(image: object, output_format: str) -> bytes:
    buffer = io.BytesIO()
    if output_format == "jpeg":
        if image.mode not in {"RGB", "L"}:
            image = image.convert("RGB")
        image.save(buffer, "JPEG", quality=85, optimize=True)
    elif output_format == "png":
        if image.mode == "P":
            image = image.convert("RGBA")
        image.save(buffer, "PNG", optimize=True)
    else:
        if image.mode == "P":
            image = image.convert("RGBA")
        image.save(buffer, "WEBP", quality=85)
    return buffer.getvalue()


def _reader(path: Path, max_input_bytes: int):
    from pypdf import PdfReader

    _check_input_size(path, max_input_bytes)
    return PdfReader(str(path), strict=True)


def _safe_metadata(reader: object) -> dict[str, str]:
    raw_metadata = getattr(reader, "metadata", None) or {}
    allowed_keys = {
        "/Title": "title",
        "/Author": "author",
        "/Subject": "subject",
        "/Creator": "creator",
        "/Producer": "producer",
        "/CreationDate": "creation_date",
        "/ModDate": "modification_date",
    }
    metadata: dict[str, str] = {}
    for raw_key, output_key in allowed_keys.items():
        value = raw_metadata.get(raw_key) if hasattr(raw_metadata, "get") else None
        if value is not None:
            metadata[output_key] = str(value)[:500]
    return metadata


def _parse_pages(raw_json: str, page_count: int, *, require_non_empty: bool) -> list[int]:
    parsed = json.loads(raw_json)
    if not isinstance(parsed, list):
        raise ValueError("pages-json must decode to a list.")
    selected: list[int] = []
    for item in parsed:
        if isinstance(item, bool):
            raise ValueError("PDF pages are 1-based positive numbers.")
        if isinstance(item, int):
            selected.append(item)
            continue
        if isinstance(item, str) and re.fullmatch(r"\d+-\d+", item):
            start_text, end_text = item.split("-", 1)
            start = int(start_text)
            end = int(end_text)
            selected.extend(range(start, end + 1))
            continue
        raise ValueError("PDF pages must be 1-based positive numbers or ranges.")
    if require_non_empty and not selected:
        raise ValueError("At least one PDF page must be selected.")
    for page_number in selected:
        if page_number < 1:
            raise ValueError("PDF page 0 is invalid; use 1-based page numbers.")
        if page_number > page_count:
            raise ValueError(f"PDF page {page_number} is out of range for {page_count} pages.")
    return selected


def _pdf_writer():
    from pypdf import PdfWriter

    return PdfWriter()


def _write_pdf(writer: object, output: Path) -> None:
    if output.exists():
        raise ValueError("Output PDF already exists.")
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as handle:
        writer.write(handle)


def _inspect_pdf(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input)
    reader = _reader(input_path, args.max_input_bytes)
    encrypted = bool(getattr(reader, "is_encrypted", False))
    page_count = 0 if encrypted else len(reader.pages)
    _write_json(
        {
            "ok": True,
            "operation": "inspect_pdf",
            "input_path": str(input_path),
            "file_size": input_path.stat().st_size,
            "page_count": page_count,
            "encrypted": encrypted,
            "metadata": _safe_metadata(reader),
        }
    )
    return 0


def _merge_pdf(args: argparse.Namespace) -> int:
    input_values = json.loads(args.inputs_json)
    if not isinstance(input_values, list) or len(input_values) < 2:
        raise ValueError("merge-pdf requires at least two input PDFs.")
    output = _validate_output_path(args.output, ".pdf")
    writer = _pdf_writer()
    total_pages = 0
    for raw_input in input_values:
        if not isinstance(raw_input, str):
            raise ValueError("merge-pdf inputs must be paths.")
        input_path = _validate_input_path(raw_input)
        reader = _reader(input_path, args.max_input_bytes)
        if reader.is_encrypted:
            raise ValueError("Encrypted PDFs are not supported.")
        for page in reader.pages:
            writer.add_page(page)
            total_pages += 1
            if total_pages > args.max_pages:
                raise ValueError(f"Merged PDF page count exceeds {args.max_pages}.")
    _write_pdf(writer, output)
    _write_json(
        {
            "ok": True,
            "operation": "merge_pdf",
            "input_file_count": len(input_values),
            "input_paths": input_values,
            "output_path": str(output),
            "output_file_size": output.stat().st_size,
            "total_pages": total_pages,
        }
    )
    return 0


def _split_pdf(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input)
    reader = _reader(input_path, args.max_input_bytes)
    if reader.is_encrypted:
        raise ValueError("Encrypted PDFs are not supported.")
    page_count = len(reader.pages)
    if page_count > args.max_output_files:
        raise ValueError(f"split-pdf would create {page_count} files; limit is {args.max_output_files}.")
    output_dir = _validate_output_dir(args.output_dir)
    stem = re.sub(r"[^A-Za-z0-9_.-]+", "_", input_path.stem).strip("._") or "document"
    created_files: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        output = output_dir / f"{stem}_page_{index:03d}.pdf"
        writer = _pdf_writer()
        writer.add_page(page)
        _write_pdf(writer, output)
        created_files.append(str(output))
    _write_json(
        {
            "ok": True,
            "operation": "split_pdf",
            "input_path": str(input_path),
            "output_dir": str(output_dir),
            "page_count": page_count,
            "created_files": created_files,
        }
    )
    return 0


def _extract_pdf(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input)
    output = _validate_output_path(args.output, ".pdf")
    reader = _reader(input_path, args.max_input_bytes)
    if reader.is_encrypted:
        raise ValueError("Encrypted PDFs are not supported.")
    selected_pages = _parse_pages(args.pages_json, len(reader.pages), require_non_empty=True)
    if len(selected_pages) > args.max_pages:
        raise ValueError(f"extract-pdf output page count exceeds {args.max_pages}.")
    writer = _pdf_writer()
    for page_number in selected_pages:
        writer.add_page(reader.pages[page_number - 1])
    _write_pdf(writer, output)
    _write_json(
        {
            "ok": True,
            "operation": "extract_pdf",
            "input_path": str(input_path),
            "output_path": str(output),
            "selected_pages": selected_pages,
            "output_page_count": len(selected_pages),
            "output_file_size": output.stat().st_size,
        }
    )
    return 0


def _rotate_pdf(args: argparse.Namespace) -> int:
    if args.rotation not in {90, 180, 270}:
        raise ValueError("Rotation must be 90, 180 or 270.")
    input_path = _validate_input_path(args.input)
    output = _validate_output_path(args.output, ".pdf")
    reader = _reader(input_path, args.max_input_bytes)
    if reader.is_encrypted:
        raise ValueError("Encrypted PDFs are not supported.")
    selected_pages = _parse_pages(args.pages_json, len(reader.pages), require_non_empty=False)
    if selected_pages:
        if len(selected_pages) > args.max_selected_pages:
            raise ValueError(f"rotate-pdf affected page count exceeds {args.max_selected_pages}.")
        selected_set = set(selected_pages)
    else:
        if len(reader.pages) > args.max_all_pages:
            raise ValueError(f"rotate-pdf all-pages output scope exceeds {args.max_all_pages} pages.")
        selected_set = set(range(1, len(reader.pages) + 1))
    writer = _pdf_writer()
    affected_pages: list[int] = []
    for index, page in enumerate(reader.pages, start=1):
        if index in selected_set:
            page.rotate(args.rotation)
            affected_pages.append(index)
        writer.add_page(page)
    _write_pdf(writer, output)
    _write_json(
        {
            "ok": True,
            "operation": "rotate_pdf",
            "input_path": str(input_path),
            "output_path": str(output),
            "rotation": args.rotation,
            "affected_pages": affected_pages,
            "output_page_count": len(reader.pages),
            "output_file_size": output.stat().st_size,
        }
    )
    return 0


def _extract_pdf_text(args: argparse.Namespace) -> int:
    if args.max_chars < 1:
        raise ValueError("max-chars must be a positive integer.")
    if args.max_preview_chars < 1:
        raise ValueError("max-preview-chars must be a positive integer.")
    input_path = _validate_input_path(args.input)
    output = None
    if args.output:
        expected_suffix = ".txt" if args.format == "text" else ".json"
        output = _validate_output_path(args.output, expected_suffix)
        if output.exists():
            raise ValueError("Output text file already exists.")
    reader = _reader(input_path, args.max_input_bytes)
    if reader.is_encrypted:
        raise ValueError("Encrypted PDFs are not supported.")
    selected_pages = _parse_pages(args.pages_json, len(reader.pages), require_non_empty=False)
    page_scope_truncated = False
    if not selected_pages:
        page_count = len(reader.pages)
        selected_pages = list(range(1, min(page_count, args.max_pages) + 1))
        page_scope_truncated = page_count > args.max_pages
    if len(selected_pages) > args.max_pages:
        raise ValueError(f"PDF page selection exceeds {args.max_pages} pages.")

    page_texts: list[dict[str, object]] = []
    char_count = 0
    truncated = False
    for page_number in selected_pages:
        text = reader.pages[page_number - 1].extract_text() or ""
        if char_count + len(text) > args.max_chars:
            remaining = max(args.max_chars - char_count, 0)
            text = text[:remaining]
            truncated = True
        page_texts.append({"page_number": page_number, "text": text})
        char_count += len(text)
        if truncated:
            break

    combined_text = "\n\n".join(str(item["text"]) for item in page_texts)
    no_text_found = not combined_text.strip()
    payload: dict[str, object] = {
        "ok": True,
        "operation": "extract_pdf_text",
        "input_path": str(input_path),
        "file_size": input_path.stat().st_size,
        "page_count": len(reader.pages),
        "extracted_pages": [item["page_number"] for item in page_texts],
        "char_count": char_count,
        "truncated": truncated,
        "page_scope_truncated": page_scope_truncated,
        "no_text_found": no_text_found,
        "format": args.format,
    }
    if args.include_preview:
        payload["preview"] = combined_text[: args.max_preview_chars]
    if output is not None:
        output.parent.mkdir(parents=True, exist_ok=True)
        if args.format == "json":
            output_payload = {
                "input_path": str(input_path),
                "page_count": len(reader.pages),
                "extracted_pages": payload["extracted_pages"],
                "char_count": char_count,
                "truncated": truncated,
                "no_text_found": no_text_found,
                "pages": page_texts,
            }
            output.write_text(json.dumps(output_payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        else:
            output.write_text(combined_text, encoding="utf-8")
        payload["output_path"] = str(output)
        payload["output_file_size"] = output.stat().st_size
    _write_json(payload)
    return 0


def _is_zip_symlink(info: zipfile.ZipInfo) -> bool:
    mode = (info.external_attr >> 16) & 0o170000
    return mode == 0o120000


def _normalize_zip_member_name(raw_name: str, *, max_path_length: int, max_depth: int) -> str:
    if "\x00" in raw_name:
        raise ValueError("ZIP entry path contains a null byte.")
    if raw_name.startswith("\\\\") or raw_name.startswith("//"):
        raise ValueError("ZIP entry uses a UNC path.")
    if re.match(r"^[A-Za-z]:[\\/]", raw_name):
        raise ValueError("ZIP entry uses a drive-qualified absolute path.")
    normalized_input = raw_name.replace("\\", "/")
    if normalized_input.startswith("/"):
        raise ValueError("ZIP entry uses an absolute path.")
    normalized = posixpath.normpath(normalized_input).strip("/")
    if normalized in {"", "."}:
        raise ValueError("ZIP entry path is empty.")
    parts = normalized.split("/")
    if any(part in {"", ".", ".."} for part in parts):
        raise ValueError("ZIP entry contains parent traversal.")
    if len(normalized) > max_path_length:
        raise ValueError(f"ZIP entry path length exceeds {max_path_length}.")
    if len(parts) > max_depth:
        raise ValueError(f"ZIP entry nesting depth exceeds {max_depth}.")
    return normalized


def _zip_entry_ratio(info: zipfile.ZipInfo) -> float:
    if info.file_size <= 0:
        return 0.0
    if info.compress_size <= 0:
        return float("inf")
    return info.file_size / max(info.compress_size, 1)


def _zip_limits_from_args(args: argparse.Namespace) -> dict[str, int]:
    return {
        "max_entries": args.max_entries,
        "max_total_uncompressed_bytes": args.max_total_uncompressed_bytes,
        "max_single_file_bytes": args.max_single_file_bytes,
        "max_path_length": args.max_path_length,
        "max_compression_ratio": args.max_compression_ratio,
        "max_depth": args.max_depth,
    }


def _inspect_zip_entries(path: Path, limits: dict[str, int]) -> dict[str, object]:
    encrypted_entries: list[str] = []
    symlink_entries: list[str] = []
    suspicious_entries: list[dict[str, object]] = []
    path_findings: list[dict[str, str]] = []
    file_entries: list[dict[str, object]] = []
    directory_count = 0
    total_compressed = 0
    total_uncompressed = 0
    largest_entry = {"name": "", "uncompressed_size": 0, "compressed_size": 0}

    with zipfile.ZipFile(path) as archive:
        infos = archive.infolist()
        if len(infos) > limits["max_entries"]:
            suspicious_entries.append({"name": "", "reason": f"entry count exceeds {limits['max_entries']}"})
        for info in infos:
            raw_name = info.filename
            try:
                normalized = _normalize_zip_member_name(
                    raw_name,
                    max_path_length=limits["max_path_length"],
                    max_depth=limits["max_depth"],
                )
            except ValueError as exc:
                path_findings.append({"entry": raw_name, "reason": str(exc)})
                normalized = raw_name
            encrypted = bool(info.flag_bits & 0x1)
            symlink = _is_zip_symlink(info)
            if encrypted:
                encrypted_entries.append(raw_name)
            if symlink:
                symlink_entries.append(raw_name)
            total_compressed += info.compress_size
            total_uncompressed += info.file_size
            if info.file_size > int(largest_entry["uncompressed_size"]):
                largest_entry = {
                    "name": normalized,
                    "uncompressed_size": info.file_size,
                    "compressed_size": info.compress_size,
                }
            ratio = _zip_entry_ratio(info)
            if info.file_size > limits["max_single_file_bytes"]:
                suspicious_entries.append({"name": raw_name, "reason": "single file size limit exceeded"})
            if info.file_size >= 1_000_000 and ratio > limits["max_compression_ratio"]:
                suspicious_entries.append({"name": raw_name, "reason": "compression ratio limit exceeded"})
            if info.is_dir():
                directory_count += 1
            else:
                file_entries.append(
                    {
                        "name": normalized,
                        "compressed_size": info.compress_size,
                        "uncompressed_size": info.file_size,
                        "compression_ratio": ratio if ratio != float("inf") else "infinite",
                    }
                )
        if total_uncompressed > limits["max_total_uncompressed_bytes"]:
            suspicious_entries.append({"name": "", "reason": "total uncompressed size limit exceeded"})

    compression_ratio = 0.0
    if total_uncompressed and total_compressed:
        compression_ratio = total_uncompressed / total_compressed
    extraction_allowed = not encrypted_entries and not symlink_entries and not suspicious_entries and not path_findings
    return {
        "entry_count": len(file_entries) + directory_count,
        "file_entry_count": len(file_entries),
        "directory_entry_count": directory_count,
        "total_compressed_bytes": total_compressed,
        "total_uncompressed_bytes": total_uncompressed,
        "largest_entry": largest_entry,
        "compression_ratio": compression_ratio,
        "encrypted_entries": encrypted_entries,
        "symlink_entries": symlink_entries,
        "suspicious_entries": suspicious_entries,
        "path_validation_findings": path_findings,
        "extraction_allowed": extraction_allowed,
        "limits": limits,
        "file_entries": file_entries,
    }


def _validate_zip_for_extraction(path: Path, limits: dict[str, int]) -> list[tuple[zipfile.ZipInfo, str]]:
    summary = _inspect_zip_entries(path, limits)
    if summary["path_validation_findings"]:
        first = summary["path_validation_findings"][0]
        raise ValueError(f"ZIP Slip/path validation rejected {first['entry']}: {first['reason']}")
    if summary["encrypted_entries"]:
        raise ValueError("Encrypted ZIP entries are not supported.")
    if summary["symlink_entries"]:
        raise ValueError("ZIP symlink entries are not supported in the MVP.")
    if summary["suspicious_entries"]:
        first = summary["suspicious_entries"][0]
        raise ValueError(f"ZIP extraction limit is exceeded: {first['reason']}")
    entries: list[tuple[zipfile.ZipInfo, str]] = []
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            normalized = _normalize_zip_member_name(
                info.filename,
                max_path_length=limits["max_path_length"],
                max_depth=limits["max_depth"],
            )
            entries.append((info, normalized))
    return entries


def _inspect_zip(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input, ".zip")
    _check_input_size(input_path, args.max_input_bytes)
    summary = _inspect_zip_entries(input_path, _zip_limits_from_args(args))
    summary.update(
        {
            "ok": True,
            "operation": "inspect_zip",
            "input_path": str(input_path),
            "file_size": input_path.stat().st_size,
        }
    )
    _write_json(summary)
    return 0


def _extract_zip(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input, ".zip")
    _check_input_size(input_path, args.max_input_bytes)
    output_dir = _validate_zip_output_dir(args.output_dir)
    limits = _zip_limits_from_args(args)
    entries = _validate_zip_for_extraction(input_path, limits)
    if output_dir.exists() and not output_dir.is_dir():
        raise ValueError("ZIP extraction output_dir exists and is not a directory.")
    output_parent = output_dir.parent
    output_parent.mkdir(parents=True, exist_ok=True)
    temp_dir = Path(tempfile.mkdtemp(prefix=f".{output_dir.name}.plwc_zip_", dir=output_parent))
    created_files: list[str] = []
    try:
        with zipfile.ZipFile(input_path) as archive:
            for info, normalized in entries:
                target = temp_dir / normalized
                resolved_target = target.resolve(strict=False)
                if resolved_target != temp_dir and temp_dir not in resolved_target.parents:
                    raise ValueError("ZIP Slip/path traversal detected during extraction.")
                if info.is_dir():
                    target.mkdir(parents=True, exist_ok=True)
                    continue
                target.parent.mkdir(parents=True, exist_ok=True)
                with archive.open(info) as source, target.open("wb") as destination:
                    shutil.copyfileobj(source, destination)
                created_files.append(str(output_dir / normalized))
        if not output_dir.exists():
            # New target dir: atomic rename of the temp staging area.
            temp_dir.replace(output_dir)
        else:
            # Existing target dir: merge after per-file collision check.
            # If anything in the staged tree would overwrite an existing
            # file in the target, fail loud before mutating output_dir.
            conflicts: list[str] = []
            for staged in sorted(temp_dir.rglob("*")):
                if not staged.is_file():
                    continue
                rel = staged.relative_to(temp_dir)
                dest = output_dir / rel
                if dest.exists():
                    conflicts.append(rel.as_posix())
            if conflicts:
                shown = ", ".join(conflicts[:5])
                more = f" (and {len(conflicts) - 5} more)" if len(conflicts) > 5 else ""
                raise ValueError(
                    f"destination_file_exists: ZIP extraction would overwrite existing files: {shown}{more}"
                )
            for staged in sorted(temp_dir.rglob("*")):
                rel = staged.relative_to(temp_dir)
                dest = output_dir / rel
                if staged.is_dir():
                    dest.mkdir(parents=True, exist_ok=True)
                    continue
                dest.parent.mkdir(parents=True, exist_ok=True)
                staged.replace(dest)
            shutil.rmtree(temp_dir, ignore_errors=True)
    except Exception:
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise
    _write_json(
        {
            "ok": True,
            "operation": "extract_zip",
            "input_path": str(input_path),
            "output_dir": str(output_dir),
            "created_files": created_files,
            "created_file_count": len(created_files),
            "limits": limits,
        }
    )
    return 0


def _collect_zip_sources(args: argparse.Namespace) -> tuple[list[tuple[Path, str]], int]:
    raw_inputs = json.loads(args.inputs_json)
    if not isinstance(raw_inputs, list) or not raw_inputs or not all(isinstance(item, str) for item in raw_inputs):
        raise ValueError("create-zip inputs-json must decode to a non-empty list of paths.")
    collected: list[tuple[Path, str]] = []
    total_size = 0
    seen: set[Path] = set()
    for raw_input in raw_inputs:
        source = _validate_zip_source_path(raw_input)
        if source.is_file():
            candidates = [source]
        elif source.is_dir():
            candidates = [path for path in source.rglob("*") if path.is_file() or path.is_symlink()]
        else:
            raise ValueError("ZIP source must be a file or directory.")
        for candidate in candidates:
            if candidate in seen:
                continue
            seen.add(candidate)
            if candidate.is_symlink():
                raise ValueError("ZIP symlink sources are not supported in the MVP.")
            size = candidate.stat().st_size
            if size > args.max_single_file_bytes:
                raise ValueError(f"ZIP source file exceeds single file limit of {args.max_single_file_bytes}.")
            total_size += size
            if total_size > args.max_total_input_bytes:
                raise ValueError(f"ZIP source total size exceeds limit of {args.max_total_input_bytes}.")
            if len(collected) + 1 > args.max_files:
                raise ValueError(f"ZIP source file count limit is {args.max_files}.")
            arcname = candidate.relative_to(WORK_ROOT.resolve(strict=True)).as_posix()
            _normalize_zip_member_name(
                arcname,
                max_path_length=args.max_path_length,
                max_depth=args.max_depth,
            )
            collected.append((candidate, arcname))
    return collected, total_size


def _create_zip(args: argparse.Namespace) -> int:
    output = _validate_output_path(args.output, ".zip")
    if output.exists():
        raise ValueError("ZIP output already exists.")
    sources, total_size = _collect_zip_sources(args)
    output.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for source, arcname in sources:
            archive.write(source, arcname)
    _write_json(
        {
            "ok": True,
            "operation": "create_zip",
            "output_path": str(output),
            "output_file_size": output.stat().st_size,
            "input_file_count": len(sources),
            "total_input_bytes": total_size,
            "entries": [arcname for _source, arcname in sources],
        }
    )
    return 0


def _read_zip_part(path: Path, member_name: str, max_bytes: int) -> bytes:
    try:
        with zipfile.ZipFile(path) as archive:
            info = archive.getinfo(member_name)
            if info.file_size > max_bytes:
                raise ValueError(f"Office XML part exceeds {max_bytes} bytes.")
            return archive.read(info)
    except KeyError as exc:
        raise ValueError(f"Malformed document: missing {member_name}.") from exc
    except zipfile.BadZipFile as exc:
        raise ValueError("Malformed document: not a valid ZIP container.") from exc


def _xml_root(raw_xml: bytes) -> Element:
    from defusedxml import ElementTree as SafeElementTree

    try:
        return SafeElementTree.fromstring(raw_xml)
    except Exception as exc:
        raise ValueError("Malformed document XML.") from exc


def _office_findings(path: Path, max_xml_part_bytes: int) -> list[str]:
    findings: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            for info in archive.infolist():
                name = info.filename
                lowered = name.casefold()
                if lowered.endswith("vbaproject.bin") or lowered.endswith(".bin") and "vba" in lowered:
                    findings.append("macro payload detected but not executed")
                if lowered.endswith(".rels"):
                    if info.file_size > max_xml_part_bytes:
                        findings.append(f"relationship part too large: {name}")
                        continue
                    try:
                        root = _xml_root(archive.read(info))
                    except ValueError:
                        findings.append(f"malformed relationship part: {name}")
                        continue
                    for element in root.iter():
                        if str(element.attrib.get("TargetMode", "")).casefold() == "external":
                            findings.append("external relationship detected but not resolved")
                            break
    except zipfile.BadZipFile as exc:
        raise ValueError("Malformed document: not a valid ZIP container.") from exc
    return sorted(set(findings))


def _truncate_text(text: str, max_chars: int) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


def _write_text_or_json_output(
    output_path: str,
    output_format: str,
    text: str,
    json_payload: dict[str, object],
) -> tuple[str | None, int | None]:
    if not output_path:
        return None, None
    expected_suffix = ".txt" if output_format == "text" else ".json"
    output = _validate_output_path(output_path, expected_suffix)
    if output.exists():
        raise ValueError("Output file already exists.")
    output.parent.mkdir(parents=True, exist_ok=True)
    if output_format == "json":
        output.write_text(json.dumps(json_payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    else:
        output.write_text(text, encoding="utf-8")
    return str(output), output.stat().st_size


def _parse_selection(raw_json: str) -> list[int | str]:
    parsed = json.loads(raw_json)
    if not isinstance(parsed, list):
        raise ValueError("pages-json must decode to a list.")
    selection: list[int | str] = []
    for item in parsed:
        if isinstance(item, bool):
            raise ValueError("Page and slide selectors are 1-based positive numbers.")
        if isinstance(item, int):
            if item < 1:
                raise ValueError("Page and slide selectors are 1-based positive numbers.")
            selection.append(item)
            continue
        if isinstance(item, str):
            stripped = item.strip()
            if re.fullmatch(r"\d+-\d+", stripped):
                start_text, end_text = stripped.split("-", 1)
                start = int(start_text)
                end = int(end_text)
                if start < 1 or end < start:
                    raise ValueError("Page and slide ranges must be ascending and 1-based.")
                selection.append(stripped)
                continue
            if stripped.isdigit() and int(stripped) >= 1:
                selection.append(int(stripped))
                continue
        raise ValueError("Page and slide selectors must be 1-based positive numbers or ranges.")
    return selection


def _expand_selection(selection: list[int | str], upper_bound: int, *, label: str) -> list[int]:
    if not selection:
        return list(range(1, upper_bound + 1))
    expanded: list[int] = []
    for item in selection:
        if isinstance(item, int):
            expanded.append(item)
        else:
            start_text, end_text = item.split("-", 1)
            expanded.extend(range(int(start_text), int(end_text) + 1))
    for number in expanded:
        if number > upper_bound:
            raise ValueError(f"{label} {number} is out of range for {upper_bound} {label.lower()}s.")
    return expanded


def _inspect_office(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input, args.expected_suffix)
    _check_input_size(input_path, args.max_input_bytes)
    command = args.command
    if command == "inspect-docx":
        payload = _inspect_docx(input_path, args.max_xml_part_bytes)
    elif command == "inspect-xlsx":
        payload = _inspect_xlsx(input_path, args.max_xml_part_bytes)
    elif command == "inspect-pptx":
        payload = _inspect_pptx(input_path, args.max_xml_part_bytes)
    elif command == "inspect-odt":
        payload = _inspect_odf_text_document(input_path, "inspect_odt", args.max_xml_part_bytes)
    elif command == "inspect-ods":
        payload = _inspect_ods(input_path, args.max_xml_part_bytes)
    elif command == "inspect-odp":
        payload = _inspect_odp(input_path, args.max_xml_part_bytes)
    else:
        raise ValueError("Unsupported Office inspect command.")
    payload.update({"ok": True, "input_path": str(input_path), "file_size": input_path.stat().st_size})
    _write_json(payload)
    return 0


def _inspect_docx(path: Path, max_xml_part_bytes: int) -> dict[str, object]:
    from docx import Document

    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs]
    table_count = len(document.tables)
    table_text_chars = sum(len(cell.text) for table in document.tables for row in table.rows for cell in row.cells)
    properties = document.core_properties
    metadata = {
        "title": properties.title or "",
        "author": properties.author or "",
        "subject": properties.subject or "",
        "keywords": properties.keywords or "",
    }
    return {
        "operation": "inspect_docx",
        "file_type": "docx",
        "paragraph_count": len(paragraphs),
        "table_count": table_count,
        "text_char_count": sum(len(item) for item in paragraphs) + table_text_chars,
        "metadata": {key: value[:500] for key, value in metadata.items() if value},
        "warnings": _office_findings(path, max_xml_part_bytes),
        "extraction_allowed": True,
    }


def _inspect_xlsx(path: Path, max_xml_part_bytes: int) -> dict[str, object]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=False, keep_links=False)
    sheets: list[dict[str, object]] = []
    formula_present = False
    for sheet in workbook.worksheets:
        sheets.append({"name": sheet.title, "max_row": sheet.max_row, "max_column": sheet.max_column})
        for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row or 0, 25), values_only=False):
            for cell in row:
                if cell.data_type == "f":
                    formula_present = True
                    break
            if formula_present:
                break
    workbook.close()
    return {
        "operation": "inspect_xlsx",
        "file_type": "xlsx",
        "sheet_count": len(sheets),
        "sheet_names": [str(sheet["name"]) for sheet in sheets],
        "sheets": sheets,
        "formula_present": formula_present,
        "warnings": _office_findings(path, max_xml_part_bytes),
        "extraction_allowed": True,
    }


def _inspect_pptx(path: Path, max_xml_part_bytes: int) -> dict[str, object]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    text_shape_count = 0
    embedded_media_count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shape_count += 1
            if getattr(shape, "shape_type", None) is not None and "PICTURE" in str(shape.shape_type):
                embedded_media_count += 1
    return {
        "operation": "inspect_pptx",
        "file_type": "pptx",
        "slide_count": len(presentation.slides),
        "text_shape_count": text_shape_count,
        "embedded_media_count": embedded_media_count,
        "warnings": _office_findings(path, max_xml_part_bytes),
        "extraction_allowed": True,
    }


def _extract_office_text(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input, args.expected_suffix)
    _check_input_size(input_path, args.max_input_bytes)
    if args.command == "extract-docx-text":
        sections, document_kind = _docx_text_sections(input_path), "docx"
    elif args.command == "extract-pptx-text":
        sections, document_kind = (
            _pptx_text_sections(input_path, args.pages_json, include_notes=bool(getattr(args, "include_notes", False))),
            "pptx",
        )
    elif args.command == "extract-odt-text":
        sections, document_kind = _odf_text_sections(input_path, args.max_xml_part_bytes, "odt"), "odt"
    elif args.command == "extract-odp-text":
        sections, document_kind = _odp_text_sections(input_path, args.max_xml_part_bytes, args.pages_json), "odp"
    else:
        raise ValueError("Unsupported Office text extraction command.")
    combined = "\n\n".join(section["text"] for section in sections if str(section.get("text", "")).strip())
    truncated_text, truncated = _truncate_text(combined, args.max_chars)
    payload = {
        "ok": True,
        "operation": args.command.replace("-", "_"),
        "input_path": str(input_path),
        "file_type": document_kind,
        "section_count": len(sections),
        "char_count": len(truncated_text),
        "truncated": truncated,
        "no_text_found": not truncated_text.strip(),
        "format": args.format,
    }
    if args.include_preview:
        payload["preview"] = truncated_text[: args.max_preview_chars]
    output_payload = payload | {"sections": sections, "text": truncated_text}
    output_path, output_size = _write_text_or_json_output(args.output, args.format, truncated_text, output_payload)
    if output_path:
        payload["output_path"] = output_path
        payload["output_file_size"] = output_size
    _write_json(payload)
    return 0


def _docx_text_sections(path: Path) -> list[dict[str, object]]:
    from docx import Document

    document = Document(str(path))
    sections: list[dict[str, object]] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if text:
            sections.append({"type": "paragraph", "index": index, "text": text})
    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            rows.append("\t".join(cell.text.strip() for cell in row.cells))
        if rows:
            sections.append({"type": "table", "index": table_index, "text": "\n".join(rows)})
    return sections


def _pptx_text_sections(path: Path, pages_json: str, *, include_notes: bool = False) -> list[dict[str, object]]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    selected = _expand_selection(_parse_selection(pages_json), len(presentation.slides), label="Slide")
    if len(selected) > OFFICE_MAX_PPTX_SLIDES:
        raise ValueError(f"PPTX slide extraction exceeds {OFFICE_MAX_PPTX_SLIDES} slides.")
    sections: list[dict[str, object]] = []
    for slide_number in selected:
        slide = presentation.slides[slide_number - 1]
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        body_text = "\n".join(texts)
        section: dict[str, object] = {"type": "slide", "slide_number": slide_number, "text": body_text}
        if include_notes:
            notes_text = ""
            if getattr(slide, "has_notes_slide", False):
                try:
                    notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                except Exception:  # pragma: no cover - delegated to python-pptx
                    notes_text = ""
            section["notes"] = notes_text
            if notes_text:
                separator = "\n" if body_text else ""
                section["text"] = f"{body_text}{separator}{notes_text}"
        sections.append(section)
    return sections


def _extract_office_data(args: argparse.Namespace) -> int:
    input_path = _validate_input_path(args.input, args.expected_suffix)
    _check_input_size(input_path, args.max_input_bytes)
    if args.command == "extract-xlsx-data":
        payload = _extract_xlsx_data(input_path, args)
    elif args.command == "extract-ods-data":
        payload = _extract_ods_data(input_path, args)
    else:
        raise ValueError("Unsupported Office data extraction command.")
    output_path = None
    output_size = None
    if args.output:
        output = _validate_output_path(args.output, ".json")
        if output.exists():
            raise ValueError("Output JSON file already exists.")
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        output_path = str(output)
        output_size = output.stat().st_size
    response = {
        key: value
        for key, value in payload.items()
        if key not in {"sheets", "tables"}
    }
    if args.include_preview:
        response["preview"] = json.dumps(payload.get("sheets") or payload.get("tables") or [], ensure_ascii=False)[: args.max_preview_chars]
    if output_path:
        response["output_path"] = output_path
        response["output_file_size"] = output_size
    _write_json(response)
    return 0


def _selected_sheet_names(available: list[str], raw_sheets_json: str, max_sheets: int) -> list[str]:
    parsed = json.loads(raw_sheets_json)
    if not isinstance(parsed, list):
        raise ValueError("sheets-json must decode to a list.")
    if not parsed:
        selected = available[:max_sheets]
    else:
        selected = []
        for item in parsed:
            if isinstance(item, bool):
                raise ValueError("Sheet selectors must be names or 1-based indexes.")
            if isinstance(item, int):
                if item < 1 or item > len(available):
                    raise ValueError("Sheet index is out of range.")
                selected.append(available[item - 1])
            elif isinstance(item, str) and item in available:
                selected.append(item)
            else:
                raise ValueError("Sheet selector is invalid.")
    if len(selected) > max_sheets:
        raise ValueError(f"Sheet selection exceeds {max_sheets} sheets.")
    return selected


def _extract_xlsx_data(path: Path, args: argparse.Namespace) -> dict[str, object]:
    from openpyxl import load_workbook
    from openpyxl.utils.cell import range_boundaries

    workbook = load_workbook(path, read_only=True, data_only=False, keep_links=False)
    selected = _selected_sheet_names(workbook.sheetnames, args.sheets_json, OFFICE_MAX_XLSX_SHEETS)
    ranges = json.loads(args.ranges_json)
    if not isinstance(ranges, (dict, list)):
        raise ValueError("ranges-json must decode to an object or list.")
    sheets: list[dict[str, object]] = []
    total_cells = 0
    formula_count = 0
    try:
        for sheet_name in selected:
            sheet = workbook[sheet_name]
            range_spec = ranges.get(sheet_name) if isinstance(ranges, dict) else (ranges[0] if ranges else None)
            if range_spec:
                min_col, min_row, max_col, max_row = range_boundaries(str(range_spec))
            else:
                min_col, min_row, max_col, max_row = 1, 1, sheet.max_column or 1, min(sheet.max_row or 1, args.max_rows)
            rows: list[list[object]] = []
            for row_index, row in enumerate(sheet.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col, values_only=False), start=1):
                if row_index > args.max_rows:
                    raise ValueError(f"XLSX row extraction exceeds {args.max_rows} rows.")
                output_row: list[object] = []
                for cell in row:
                    total_cells += 1
                    if total_cells > args.max_cells:
                        raise ValueError(f"XLSX cell extraction exceeds {args.max_cells} cells.")
                    if cell.data_type == "f":
                        formula_count += 1
                        output_row.append({"type": "formula", "formula": str(cell.value)})
                    else:
                        output_row.append(cell.value)
                rows.append(output_row)
            sheets.append({"name": sheet_name, "range": range_spec or sheet.calculate_dimension(), "rows": rows})
    finally:
        workbook.close()
    return {
        "ok": True,
        "operation": "extract_xlsx_data",
        "input_path": str(path),
        "file_type": "xlsx",
        "sheet_count": len(sheets),
        "extracted_cells": total_cells,
        "formula_count": formula_count,
        "formula_policy": "Formulas are reported as formula text and are not executed.",
        "sheets": sheets,
        "format": "json",
    }


ODF_NS = {
    "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
    "draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
}


def _odf_root(path: Path, max_xml_part_bytes: int) -> Element:
    return _xml_root(_read_zip_part(path, "content.xml", max_xml_part_bytes))


def _node_text(node: Element) -> str:
    return " ".join(text.strip() for text in node.itertext() if text and text.strip())


def _inspect_odf_text_document(path: Path, operation: str, max_xml_part_bytes: int) -> dict[str, object]:
    root = _odf_root(path, max_xml_part_bytes)
    paragraphs = root.findall(".//text:p", ODF_NS)
    tables = root.findall(".//table:table", ODF_NS)
    return {
        "operation": operation,
        "file_type": path.suffix.lower().lstrip("."),
        "paragraph_count": len(paragraphs),
        "table_count": len(tables),
        "text_char_count": sum(len(_node_text(node)) for node in paragraphs),
        "warnings": _office_findings(path, max_xml_part_bytes),
        "extraction_allowed": True,
    }


def _inspect_ods(path: Path, max_xml_part_bytes: int) -> dict[str, object]:
    root = _odf_root(path, max_xml_part_bytes)
    tables = root.findall(".//table:table", ODF_NS)
    table_payload: list[dict[str, object]] = []
    formula_present = False
    for table in tables:
        rows = table.findall("table:table-row", ODF_NS)
        cells = table.findall(".//table:table-cell", ODF_NS)
        for cell in cells:
            if any(key.endswith("formula") for key in cell.attrib):
                formula_present = True
                break
        table_payload.append({"name": table.attrib.get(f"{{{ODF_NS['table']}}}name", ""), "row_count": len(rows), "cell_count": len(cells)})
    return {
        "operation": "inspect_ods",
        "file_type": "ods",
        "table_count": len(tables),
        "sheet_names": [str(table["name"]) for table in table_payload],
        "tables": table_payload,
        "formula_present": formula_present,
        "warnings": _office_findings(path, max_xml_part_bytes),
        "extraction_allowed": True,
    }


def _inspect_odp(path: Path, max_xml_part_bytes: int) -> dict[str, object]:
    root = _odf_root(path, max_xml_part_bytes)
    pages = root.findall(".//draw:page", ODF_NS)
    text_elements = root.findall(".//text:p", ODF_NS)
    return {
        "operation": "inspect_odp",
        "file_type": "odp",
        "slide_count": len(pages),
        "text_element_count": len(text_elements),
        "warnings": _office_findings(path, max_xml_part_bytes),
        "extraction_allowed": True,
    }


def _odf_text_sections(path: Path, max_xml_part_bytes: int, file_type: str) -> list[dict[str, object]]:
    root = _odf_root(path, max_xml_part_bytes)
    sections: list[dict[str, object]] = []
    for index, paragraph in enumerate(root.findall(".//text:p", ODF_NS), start=1):
        text = _node_text(paragraph)
        if text:
            sections.append({"type": "paragraph", "index": index, "text": text})
    for table_index, table in enumerate(root.findall(".//table:table", ODF_NS), start=1):
        rows: list[str] = []
        for row in table.findall("table:table-row", ODF_NS):
            cells = [_node_text(cell) for cell in row.findall("table:table-cell", ODF_NS)]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            sections.append({"type": "table", "index": table_index, "text": "\n".join(rows)})
    return sections


def _odp_text_sections(path: Path, max_xml_part_bytes: int, pages_json: str) -> list[dict[str, object]]:
    root = _odf_root(path, max_xml_part_bytes)
    pages = root.findall(".//draw:page", ODF_NS)
    selected = _expand_selection(_parse_selection(pages_json), len(pages), label="Slide")
    if len(selected) > OFFICE_MAX_PPTX_SLIDES:
        raise ValueError(f"ODP slide extraction exceeds {OFFICE_MAX_PPTX_SLIDES} slides.")
    sections: list[dict[str, object]] = []
    for page_number in selected:
        page = pages[page_number - 1]
        text = "\n".join(_node_text(node) for node in page.findall(".//text:p", ODF_NS) if _node_text(node))
        sections.append({"type": "slide", "slide_number": page_number, "text": text})
    return sections


def _extract_ods_data(path: Path, args: argparse.Namespace) -> dict[str, object]:
    root = _odf_root(path, args.max_xml_part_bytes)
    tables = root.findall(".//table:table", ODF_NS)
    available = [table.attrib.get(f"{{{ODF_NS['table']}}}name", f"Table {index}") for index, table in enumerate(tables, start=1)]
    selected_names = set(_selected_sheet_names(available, args.sheets_json, OFFICE_MAX_ODF_TABLES))
    output_tables: list[dict[str, object]] = []
    total_cells = 0
    formula_count = 0
    for table, name in zip(tables, available):
        if name not in selected_names:
            continue
        rows_payload: list[list[object]] = []
        for row_index, row in enumerate(table.findall("table:table-row", ODF_NS), start=1):
            if row_index > args.max_rows:
                raise ValueError(f"ODS row extraction exceeds {args.max_rows} rows.")
            row_payload: list[object] = []
            for cell in row.findall("table:table-cell", ODF_NS):
                total_cells += 1
                if total_cells > args.max_cells:
                    raise ValueError(f"ODS cell extraction exceeds {args.max_cells} cells.")
                formula = next((value for key, value in cell.attrib.items() if key.endswith("formula")), None)
                if formula:
                    formula_count += 1
                    row_payload.append({"type": "formula", "formula": formula, "text": _node_text(cell)})
                else:
                    row_payload.append(_node_text(cell))
            if any(value not in {"", None} for value in row_payload):
                rows_payload.append(row_payload)
        output_tables.append({"name": name, "rows": rows_payload})
    return {
        "ok": True,
        "operation": "extract_ods_data",
        "input_path": str(path),
        "file_type": "ods",
        "table_count": len(output_tables),
        "extracted_cells": total_cells,
        "formula_count": formula_count,
        "formula_policy": "Formulas are reported as formula text and are not executed.",
        "tables": output_tables,
        "format": "json",
    }


def _parse_content(raw_json: str, command: str = "") -> dict[str, object]:
    parsed = json.loads(raw_json)
    if not isinstance(parsed, dict):
        raise ValueError("content-json must decode to an object.")
    # XLSX V2 stores formula-looking strings as literal text in value cells and
    # writes only explicit formula cells; PPTX V2 stores arbitrary slide and
    # notes text verbatim and never executes or interprets it; PDF V2 stores
    # arbitrary paragraph/table/list text verbatim and never interprets it.
    # All three creation paths therefore opt out of the legacy formula-prefix
    # rejection.
    _reject_active_strings(
        parsed, allow_formula_literals=command in {"create-xlsx", "create-pptx", "create-pdf"}
    )
    return parsed


def _creation_content_from_args(args: argparse.Namespace) -> tuple[dict[str, object], dict[str, object]]:
    input_json = getattr(args, "input_json", "")
    if not input_json:
        return _parse_content(args.content_json, args.command), {}
    if args.command not in {"create-docx", "create-xlsx", "create-pptx", "create-pdf"}:
        raise ValueError(
            "validation_error: only create-docx, create-xlsx, create-pptx and create-pdf support input-json."
        )
    input_path = _validate_input_path(input_json, ".json")
    input_size = _check_input_size(input_path, args.max_json_input_bytes)
    content = _parse_content(input_path.read_text(encoding="utf-8"), args.command)
    return content, {"input_path": str(input_path), "input_file_size": input_size}


def _reject_active_strings(value: object, *, allow_formula_literals: bool = False) -> None:
    if isinstance(value, str):
        if value.lstrip().startswith("=") and not allow_formula_literals:
            raise ValueError("Spreadsheet formulas and active content are not supported.")
        return
    if isinstance(value, list):
        for item in value:
            _reject_active_strings(item, allow_formula_literals=allow_formula_literals)
        return
    if isinstance(value, dict):
        for item in value.values():
            _reject_active_strings(item, allow_formula_literals=allow_formula_literals)


def _string_list(value: object, fallback: list[str]) -> list[str]:
    if isinstance(value, list) and all(isinstance(item, str) for item in value):
        return value
    return fallback


def _edit_docx(args: argparse.Namespace) -> int:
    """Governed declarative DOCX editor: read-A-write-B semantics."""
    input_path = _validate_input_path(args.input, ".docx")
    output_path = _validate_output_path(args.output, ".docx")
    _check_input_size(input_path, args.max_input_bytes)

    edits: list[dict[str, object]] = json.loads(args.edits_json)

    # Security: scan for macros and external relationships before loading with python-docx.
    findings = _office_findings(input_path, args.max_xml_part_bytes)
    macro_findings = [f for f in findings if "macro" in f]
    if macro_findings:
        _write_json({
            "ok": False,
            "operation": "edit_docx",
            "error": f"edit_docx rejected: {macro_findings[0]}",
            "error_category": "macro_detected",
        })
        return 1

    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    document = Document(str(input_path))
    ops_applied = 0

    for edit in edits:
        op_type = str(edit.get("type", ""))

        if op_type == "replace_text":
            find = str(edit["find"])
            replace = str(edit["replace"])
            max_replacements = edit.get("max_replacements")
            limit: int | None = int(max_replacements) if max_replacements is not None else None
            _edit_replace_text(document, find, replace, limit)
            ops_applied += 1

        elif op_type == "set_core_property":
            name = str(edit["name"])
            value = str(edit["value"])
            props = document.core_properties
            if name == "title":
                props.title = value
            elif name == "author":
                props.author = value
            elif name == "subject":
                props.subject = value
            elif name == "keywords":
                props.keywords = value
            ops_applied += 1

        elif op_type == "append_paragraph":
            text = str(edit["text"])
            style_key = edit.get("style")
            para = document.add_paragraph(text)
            if style_key:
                docx_style = EDIT_DOCX_STYLE_MAP.get(str(style_key), "Normal")
                try:
                    para.style = docx_style
                except KeyError:
                    pass  # style not in document template; leave as Normal
            ops_applied += 1

        elif op_type in ("set_footer_text", "set_header_text"):
            text = str(edit["text"])
            add_page_num = bool(edit.get("page_number", False))
            section = document.sections[0]
            hdr_or_ftr = section.header if op_type == "set_header_text" else section.footer
            hdr_or_ftr.is_linked_to_previous = False
            # Clear existing paragraph content
            para = hdr_or_ftr.paragraphs[0] if hdr_or_ftr.paragraphs else hdr_or_ftr.add_paragraph()
            p_elem = para._p
            for r_elem in p_elem.findall(qn("w:r")):
                p_elem.remove(r_elem)
            for hyperlink in p_elem.findall(qn("w:hyperlink")):
                p_elem.remove(hyperlink)
            # Add text run
            if text:
                text_run = para.add_run(text)  # noqa: F841
            # Optionally append PAGE field
            if add_page_num:
                field_run = para.add_run()
                fld_begin = OxmlElement("w:fldChar")
                fld_begin.set(qn("w:fldCharType"), "begin")
                field_run._r.append(fld_begin)
                instr = OxmlElement("w:instrText")
                instr.text = " PAGE "
                field_run._r.append(instr)
                fld_sep = OxmlElement("w:fldChar")
                fld_sep.set(qn("w:fldCharType"), "separate")
                field_run._r.append(fld_sep)
                fld_end = OxmlElement("w:fldChar")
                fld_end.set(qn("w:fldCharType"), "end")
                field_run._r.append(fld_end)
            ops_applied += 1

    document.save(str(output_path))

    _write_json({
        "ok": True,
        "operation": "edit_docx",
        "output_path": str(output_path),
        "file_size": output_path.stat().st_size,
        "ops_applied": ops_applied,
    })
    return 0


def _iter_all_paragraphs(document: object):  # type: ignore[no-untyped-def]
    """Yield all paragraphs: body + table cells."""
    from docx import Document as _Doc  # noqa: F401
    yield from document.paragraphs
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs


def _replace_in_string(text: str, find: str, replace: str, remaining: int | None) -> tuple[str, int]:
    if remaining is None:
        return text.replace(find, replace), text.count(find)
    result: list[str] = []
    count = 0
    start = 0
    while start <= len(text):
        pos = text.find(find, start)
        if pos == -1:
            result.append(text[start:])
            break
        if count >= remaining:
            result.append(text[start:])
            break
        result.append(text[start:pos])
        result.append(replace)
        count += 1
        start = pos + len(find)
    return "".join(result), count


def _edit_replace_text(document: object, find: str, replace: str, limit: int | None) -> int:  # type: ignore[no-untyped-def]
    """Replace occurrences of `find` with `replace` across all body/table paragraphs."""
    total = 0
    scanned = 0
    for para in _iter_all_paragraphs(document):
        if scanned >= EDIT_DOCX_MAX_SCANNED_PARAGRAPHS:
            break
        scanned += 1
        for run in para.runs:
            if find not in run.text:
                continue
            remaining = (limit - total) if limit is not None else None
            if remaining is not None and remaining <= 0:
                return total
            new_text, n = _replace_in_string(run.text, find, replace, remaining)
            run.text = new_text
            total += n
    return total


def _create_docx(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    if _is_docx_v2_content(content):
        return _create_docx_v2(output, title, body, content)
    return _create_docx_legacy(output, title, body, content)


def _is_docx_v2_content(content: dict[str, object]) -> bool:
    return any(key in content for key in ("schema_version", "document", "page", "styles", "content"))


def _create_docx_legacy(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    from docx import Document

    document_title = str(content.get("title") or title)
    paragraphs = _string_list(content.get("paragraphs"), [body])
    table_rows = content.get("table")
    document = Document()
    document.add_heading(document_title, level=1)
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    if isinstance(table_rows, list) and table_rows and all(isinstance(row, list) for row in table_rows):
        max_cols = max(len(row) for row in table_rows if isinstance(row, list))
        if max_cols:
            table = document.add_table(rows=len(table_rows), cols=max_cols)
            for row_index, row in enumerate(table_rows):
                if not isinstance(row, list):
                    continue
                for col_index, cell in enumerate(row[:max_cols]):
                    table.cell(row_index, col_index).text = "" if cell is None else str(cell)
    else:
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Field"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Worker"
        table.cell(1, 1).text = "PLwC Document Worker"
    document.save(output)
    return {
        "builder_version": "docx_legacy_compat",
        "document_format": "docx",
        "content_element_count": len(paragraphs) + 1,
        "table_count": 1,
        "image_count": 0,
    }


def _create_docx_v2(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    from docx import Document
    from docx.enum.section import WD_ORIENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Mm, Pt

    _validate_docx_v2_worker_content(content)
    document = Document()
    metadata = content.get("document") if isinstance(content.get("document"), dict) else {}
    document_title = str(metadata.get("title") or title)
    if document_title:
        document.core_properties.title = document_title
    if isinstance(metadata.get("author"), str):
        document.core_properties.author = str(metadata["author"])

    page_settings = _docx_v2_page_settings(content.get("page"))
    section = document.sections[0]
    width_mm, height_mm = DOCX_V2_PAGE_SIZES_MM[page_settings["size"]]
    if page_settings["orientation"] == "landscape":
        width_mm, height_mm = height_mm, width_mm
        section.orientation = WD_ORIENT.LANDSCAPE
    else:
        section.orientation = WD_ORIENT.PORTRAIT
    section.page_width = Mm(width_mm)
    section.page_height = Mm(height_mm)
    margins = page_settings["margins_mm"]
    section.top_margin = Mm(margins["top"])
    section.bottom_margin = Mm(margins["bottom"])
    section.left_margin = Mm(margins["left"])
    section.right_margin = Mm(margins["right"])

    style_specs = content.get("styles") if isinstance(content.get("styles"), dict) else {}
    _apply_docx_v2_styles(document, style_specs)

    alignments = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
    }
    image_alignments = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
    }
    elements = content["content"]
    paragraph_count = 0
    table_count = 0
    image_count = 0
    for element in elements:
        element_type = str(element["type"])
        if element_type == "page_break":
            document.add_page_break()
            continue
        if element_type in DOCX_V2_TEXT_TYPES:
            paragraph_count += 1
            paragraph = _add_docx_v2_text_paragraph(document, element)
            override = element.get("style_overrides")
            if isinstance(override, dict):
                _apply_docx_v2_paragraph_spec(paragraph, override, alignments)
                _apply_docx_v2_run_spec(paragraph.runs, override)
            continue
        if element_type in {"bullet_list", "numbered_list"}:
            paragraph_count += _add_docx_v2_list(document, element)
            continue
        if element_type == "table":
            table_count += 1
            _add_docx_v2_table(document, element)
            continue
        if element_type == "image":
            image_count += 1
            paragraph = document.add_paragraph()
            paragraph.alignment = image_alignments.get(str(element.get("align", "left")).casefold(), WD_ALIGN_PARAGRAPH.LEFT)
            image_path = _resolve_docx_v2_image_path(str(element["path"]))
            width = Mm(float(element["width_mm"])) if element.get("width_mm") is not None else None
            height = Mm(float(element["height_mm"])) if element.get("height_mm") is not None else None
            run = paragraph.add_run()
            if width is not None and height is not None:
                run.add_picture(str(image_path), width=width, height=height)
            elif width is not None:
                run.add_picture(str(image_path), width=width)
            elif height is not None:
                run.add_picture(str(image_path), height=height)
            else:
                run.add_picture(str(image_path))

    document.save(output)
    return {
        "builder_version": "docx_creation_v2",
        "document_format": "docx",
        "content_element_count": len(elements),
        "paragraph_count": paragraph_count,
        "table_count": table_count,
        "image_count": image_count,
        "page_size": page_settings["size"],
        "orientation": page_settings["orientation"],
        "limits": _docx_v2_limits(),
    }


def _validate_docx_v2_worker_content(content: dict[str, object]) -> None:
    unknown_top = sorted(str(key) for key in content if str(key) not in {"schema_version", "document", "page", "styles", "content"})
    if unknown_top:
        raise ValueError(f"validation_error: create_docx unsupported fields: {', '.join(unknown_top)}.")
    schema_version = content.get("schema_version", DOCX_V2_SCHEMA_VERSION)
    if not isinstance(schema_version, str) or schema_version not in DOCX_V2_SUPPORTED_SCHEMA_VERSIONS:
        raise ValueError(f"validation_error: create_docx schema_version must be one of: {', '.join(sorted(DOCX_V2_SUPPORTED_SCHEMA_VERSIONS))}.")
    elements = content.get("content")
    if not isinstance(elements, list) or not elements:
        raise ValueError("validation_error: create_docx content must be a non-empty list.")
    if len(elements) > DOCX_V2_MAX_CONTENT_ELEMENTS:
        raise ValueError(f"create_docx content element count exceeds {DOCX_V2_MAX_CONTENT_ELEMENTS}.")
    paragraphs = 0
    tables = 0
    cells = 0
    images = 0
    for index, element in enumerate(elements, start=1):
        if not isinstance(element, dict):
            raise ValueError(f"validation_error: create_docx element {index} must be an object.")
        element_type = element.get("type")
        if not isinstance(element_type, str) or element_type not in DOCX_V2_CONTENT_TYPES:
            raise ValueError(f"validation_error: create_docx element {index} has an unsupported type.")
        if element_type in DOCX_V2_TEXT_TYPES:
            paragraphs += 1
            _validate_docx_v2_worker_text_element(index, element_type, element)
        elif element_type in {"bullet_list", "numbered_list"}:
            items = _validate_docx_v2_worker_list_element(index, element_type, element)
            paragraphs += len(items)
        elif element_type == "page_break":
            unknown = sorted(str(key) for key in element if str(key) != "type")
            if unknown:
                raise ValueError(f"validation_error: create_docx page_break element {index} has unsupported fields: {', '.join(unknown)}.")
        elif element_type == "table":
            tables += 1
            rows = element.get("rows")
            if not isinstance(rows, list) or not rows:
                raise ValueError(f"validation_error: create_docx table element {index}.rows must be a non-empty list.")
            for row in rows:
                if not isinstance(row, list):
                    raise ValueError(f"validation_error: create_docx table element {index} rows must be arrays.")
                cells += len(row)
        elif element_type == "image":
            images += 1
            _validate_docx_v2_worker_image_element(index, element)
    if paragraphs > DOCX_V2_MAX_PARAGRAPHS:
        raise ValueError(f"create_docx paragraph-like element count exceeds {DOCX_V2_MAX_PARAGRAPHS}.")
    if tables > DOCX_V2_MAX_TABLES:
        raise ValueError(f"create_docx table count exceeds {DOCX_V2_MAX_TABLES}.")
    if cells > DOCX_V2_MAX_TABLE_CELLS:
        raise ValueError(f"create_docx table cell count exceeds {DOCX_V2_MAX_TABLE_CELLS}.")
    if images > DOCX_V2_MAX_IMAGES:
        raise ValueError(f"create_docx image count exceeds {DOCX_V2_MAX_IMAGES}.")


def _validate_docx_v2_worker_text_element(index: int, element_type: str, element: dict[str, object]) -> None:
    unknown = sorted(str(key) for key in element if str(key) not in {"type", "text", "runs", "style", "style_overrides"})
    if unknown:
        raise ValueError(f"validation_error: create_docx {element_type} element {index} has unsupported fields: {', '.join(unknown)}.")
    has_text = "text" in element
    has_runs = "runs" in element
    if has_text and has_runs:
        raise ValueError(f"validation_error: create_docx {element_type} element {index} must use either text or runs, not both.")
    if not has_text and not has_runs:
        raise ValueError(f"validation_error: create_docx {element_type} element {index} requires text or runs.")
    if has_text:
        text = element.get("text")
        if not isinstance(text, str):
            raise ValueError(f"validation_error: create_docx {element_type} element {index}.text must be a string.")
        if element_type != "paragraph" and not text.strip():
            raise ValueError(f"validation_error: create_docx {element_type} element {index}.text must not be empty.")
    if has_runs:
        runs = element.get("runs")
        if not isinstance(runs, list) or not runs:
            raise ValueError(f"validation_error: create_docx {element_type} element {index}.runs must be a non-empty list.")
        for run_index, run in enumerate(runs, start=1):
            _validate_docx_v2_worker_run(index, run_index, run)
    style_name = element.get("style")
    if style_name is not None and (not isinstance(style_name, str) or style_name not in DOCX_V2_STYLE_NAMES):
        raise ValueError(f"validation_error: create_docx {element_type} element {index}.style is not supported.")
    overrides = element.get("style_overrides")
    if overrides is not None:
        if not isinstance(overrides, dict):
            raise ValueError(f"validation_error: create_docx {element_type} element {index}.style_overrides must be an object.")
        _validate_docx_v2_style_spec(f"content[{index}].style_overrides", overrides)


def _validate_docx_v2_worker_run(element_index: int, run_index: int, run: object) -> None:
    if not isinstance(run, dict):
        raise ValueError(f"validation_error: create_docx content element {element_index} run {run_index} must be an object.")
    unknown = sorted(str(key) for key in run if str(key) not in DOCX_V2_RUN_FIELDS)
    if unknown:
        raise ValueError(f"validation_error: create_docx content element {element_index} run {run_index} has unsupported fields: {', '.join(unknown)}.")
    if not isinstance(run.get("text"), str):
        raise ValueError(f"validation_error: create_docx content element {element_index} run {run_index}.text must be a string.")
    for field_name in ("bold", "italic", "underline"):
        value = run.get(field_name)
        if value is not None and not isinstance(value, bool):
            raise ValueError(f"validation_error: create_docx content element {element_index} run {run_index}.{field_name} must be a boolean.")


def _validate_docx_v2_worker_list_element(index: int, element_type: str, element: dict[str, object]) -> list[str]:
    unknown = sorted(str(key) for key in element if str(key) not in {"type", "items", "style_overrides"})
    if unknown:
        raise ValueError(f"validation_error: create_docx {element_type} element {index} has unsupported fields: {', '.join(unknown)}.")
    items = element.get("items")
    if not isinstance(items, list) or not items:
        raise ValueError(f"validation_error: create_docx {element_type} element {index}.items must be a non-empty list of strings.")
    for item_index, item in enumerate(items, start=1):
        if not isinstance(item, str) or not item.strip():
            raise ValueError(f"validation_error: create_docx {element_type} element {index} item {item_index} must be a non-empty string.")
    overrides = element.get("style_overrides")
    if overrides is not None:
        if not isinstance(overrides, dict):
            raise ValueError(f"validation_error: create_docx {element_type} element {index}.style_overrides must be an object.")
        _validate_docx_v2_style_spec(f"content[{index}].style_overrides", overrides)
    return items


def _validate_docx_v2_worker_image_element(index: int, element: dict[str, object]) -> None:
    unknown = sorted(str(key) for key in element if str(key) not in {"type", "path", "width_mm", "height_mm", "align"})
    if unknown:
        raise ValueError(f"validation_error: create_docx image element {index} has unsupported fields: {', '.join(unknown)}.")
    if not isinstance(element.get("path"), str) or not str(element["path"]).strip():
        raise ValueError(f"validation_error: create_docx image element {index}.path must be a string.")
    align = str(element.get("align", "left")).casefold()
    if align not in {"left", "center", "right"}:
        raise ValueError("validation_error: create_docx image align must be left, center or right.")
    for field_name in ("width_mm", "height_mm"):
        value = element.get(field_name)
        if value is not None and (not isinstance(value, (int, float)) or isinstance(value, bool) or not 1 <= float(value) <= 300):
            raise ValueError(f"validation_error: create_docx image {field_name} must be a number between 1 and 300.")
    _resolve_docx_v2_image_path(str(element["path"]))


def _docx_v2_page_settings(page: object) -> dict[str, object]:
    default = {
        "size": "A4",
        "orientation": "portrait",
        "margins_mm": {"top": 18, "bottom": 18, "left": 16, "right": 16},
    }
    if page is None:
        return default
    if not isinstance(page, dict):
        raise ValueError("validation_error: create_docx page must be an object.")
    unknown = sorted(str(key) for key in page if str(key) not in {"size", "orientation", "margins_mm"})
    if unknown:
        raise ValueError(f"validation_error: create_docx page has unsupported fields: {', '.join(unknown)}.")
    size = str(page.get("size", default["size"])).upper()
    if size not in DOCX_V2_PAGE_SIZES_MM:
        raise ValueError("validation_error: create_docx page.size must be A4 or A5.")
    orientation = str(page.get("orientation", default["orientation"])).casefold()
    if orientation not in {"portrait", "landscape"}:
        raise ValueError("validation_error: create_docx page.orientation must be portrait or landscape.")
    margins = dict(default["margins_mm"])
    raw_margins = page.get("margins_mm")
    if raw_margins is not None:
        if not isinstance(raw_margins, dict):
            raise ValueError("validation_error: create_docx page.margins_mm must be an object.")
        unknown_margins = sorted(str(key) for key in raw_margins if str(key) not in {"top", "bottom", "left", "right"})
        if unknown_margins:
            raise ValueError(f"validation_error: create_docx margins has unsupported fields: {', '.join(unknown_margins)}.")
        for key, value in raw_margins.items():
            if not isinstance(value, (int, float)) or isinstance(value, bool) or not 0 <= float(value) <= 100:
                raise ValueError(f"validation_error: create_docx margin {key} must be a number between 0 and 100.")
            margins[str(key)] = float(value)
    return {"size": size, "orientation": orientation, "margins_mm": margins}


def _apply_docx_v2_styles(document: object, styles: object) -> None:
    if not isinstance(styles, dict):
        return
    unknown_styles = sorted(str(key) for key in styles if str(key) not in DOCX_V2_STYLE_NAMES)
    if unknown_styles:
        raise ValueError(f"validation_error: create_docx styles has unsupported names: {', '.join(unknown_styles)}.")
    mapping = {
        "body": "Normal",
        "title": "Title",
        "heading1": "Heading 1",
        "heading2": "Heading 2",
        "heading3": "Heading 3",
        "quote": "Intense Quote",
        "blockquote": "Intense Quote",
    }
    for public_name, spec in styles.items():
        if not isinstance(spec, dict):
            raise ValueError(f"validation_error: create_docx styles.{public_name} must be an object.")
        _validate_docx_v2_style_spec(f"styles.{public_name}", spec)
        docx_style_name = mapping[str(public_name)]
        try:
            style = document.styles[docx_style_name]
        except KeyError:
            continue
        _apply_docx_v2_style_spec(style, spec)


def _validate_docx_v2_style_spec(path: str, spec: dict[str, object]) -> None:
    unknown = sorted(str(key) for key in spec if str(key) not in DOCX_V2_STYLE_FIELDS)
    if unknown:
        raise ValueError(f"validation_error: create_docx {path} has unsupported fields: {', '.join(unknown)}.")
    for field_name, value in spec.items():
        if field_name == "font" and (not isinstance(value, str) or not value.strip() or len(value) > 100):
            raise ValueError(f"validation_error: create_docx {path}.font must be a non-empty string up to 100 characters.")
        if field_name == "size_pt" and (not isinstance(value, (int, float)) or isinstance(value, bool) or not 1 <= float(value) <= 96):
            raise ValueError(f"validation_error: create_docx {path}.size_pt must be a number between 1 and 96.")
        if field_name in {"bold", "italic", "page_break_before"} and not isinstance(value, bool):
            raise ValueError(f"validation_error: create_docx {path}.{field_name} must be a boolean.")
        if field_name == "alignment" and (not isinstance(value, str) or value.casefold() not in {"left", "center", "right", "justify"}):
            raise ValueError(f"validation_error: create_docx {path}.alignment must be left, center, right or justify.")
        if field_name == "line_spacing" and (not isinstance(value, (int, float)) or isinstance(value, bool) or not 0.5 <= float(value) <= 3):
            raise ValueError(f"validation_error: create_docx {path}.line_spacing must be a number between 0.5 and 3.")
        if field_name in {"space_before_pt", "space_after_pt"} and (
            not isinstance(value, (int, float)) or isinstance(value, bool) or not 0 <= float(value) <= 144
        ):
            raise ValueError(f"validation_error: create_docx {path}.{field_name} must be a number between 0 and 144.")
        if field_name == "first_line_indent_mm" and (
            not isinstance(value, (int, float)) or isinstance(value, bool) or not 0 <= float(value) <= 50
        ):
            raise ValueError(f"validation_error: create_docx {path}.first_line_indent_mm must be a number between 0 and 50.")


def _apply_docx_v2_style_spec(style: object, spec: dict[str, object]) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Mm, Pt

    if "font" in spec:
        style.font.name = str(spec["font"])
    if "size_pt" in spec:
        style.font.size = Pt(float(spec["size_pt"]))
    if "bold" in spec:
        style.font.bold = bool(spec["bold"])
    if "italic" in spec:
        style.font.italic = bool(spec["italic"])
    paragraph_format = style.paragraph_format
    if "line_spacing" in spec:
        paragraph_format.line_spacing = float(spec["line_spacing"])
    if "space_before_pt" in spec:
        paragraph_format.space_before = Pt(float(spec["space_before_pt"]))
    if "space_after_pt" in spec:
        paragraph_format.space_after = Pt(float(spec["space_after_pt"]))
    if "first_line_indent_mm" in spec:
        paragraph_format.first_line_indent = Mm(float(spec["first_line_indent_mm"]))
    if "page_break_before" in spec:
        paragraph_format.page_break_before = bool(spec["page_break_before"])
    if "alignment" in spec:
        paragraph_format.alignment = {
            "left": WD_ALIGN_PARAGRAPH.LEFT,
            "center": WD_ALIGN_PARAGRAPH.CENTER,
            "right": WD_ALIGN_PARAGRAPH.RIGHT,
            "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
        }[str(spec["alignment"]).casefold()]


def _add_docx_v2_text_paragraph(document: object, element: dict[str, object]) -> object:
    element_type = str(element["type"])
    style_name = str(element.get("style") or element_type)
    style_map = {
        "body": "Normal",
        "paragraph": "Normal",
        "title": "Title",
        "heading1": "Heading 1",
        "heading2": "Heading 2",
        "heading3": "Heading 3",
        "quote": "Intense Quote",
        "blockquote": "Intense Quote",
    }
    paragraph = document.add_paragraph(style=style_map.get(style_name, style_map.get(element_type, "Normal")))
    if isinstance(element.get("runs"), list):
        for run_spec in element["runs"]:
            run = paragraph.add_run(str(run_spec.get("text", "")))
            if "bold" in run_spec:
                run.bold = bool(run_spec["bold"])
            if "italic" in run_spec:
                run.italic = bool(run_spec["italic"])
            if "underline" in run_spec:
                run.underline = bool(run_spec["underline"])
    else:
        paragraph.add_run(str(element.get("text", "")))
    return paragraph


def _apply_docx_v2_paragraph_spec(paragraph: object, spec: dict[str, object], alignments: dict[str, object]) -> None:
    from docx.shared import Mm, Pt

    if "alignment" in spec:
        paragraph.alignment = alignments[str(spec["alignment"]).casefold()]
    if "line_spacing" in spec:
        paragraph.paragraph_format.line_spacing = float(spec["line_spacing"])
    if "space_before_pt" in spec:
        paragraph.paragraph_format.space_before = Pt(float(spec["space_before_pt"]))
    if "space_after_pt" in spec:
        paragraph.paragraph_format.space_after = Pt(float(spec["space_after_pt"]))
    if "first_line_indent_mm" in spec:
        paragraph.paragraph_format.first_line_indent = Mm(float(spec["first_line_indent_mm"]))
    if "page_break_before" in spec:
        paragraph.paragraph_format.page_break_before = bool(spec["page_break_before"])


def _apply_docx_v2_run_spec(runs: object, spec: dict[str, object]) -> None:
    from docx.shared import Pt

    for run in runs:
        if "font" in spec:
            run.font.name = str(spec["font"])
        if "size_pt" in spec:
            run.font.size = Pt(float(spec["size_pt"]))
        if "bold" in spec:
            run.bold = bool(spec["bold"])
        if "italic" in spec:
            run.italic = bool(spec["italic"])


def _add_docx_v2_table(document: object, element: dict[str, object]) -> None:
    rows = element["rows"]
    max_cols = max(len(row) for row in rows)
    table = document.add_table(rows=len(rows), cols=max_cols)
    try:
        table.style = "Table Grid"
    except KeyError:
        pass
    for row_index, row in enumerate(rows):
        for col_index in range(max_cols):
            value = row[col_index] if col_index < len(row) else ""
            table.cell(row_index, col_index).text = "" if value is None else str(value)


def _add_docx_v2_list(document: object, element: dict[str, object]) -> int:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    style = "List Bullet" if element["type"] == "bullet_list" else "List Number"
    items = element["items"]
    for item in items:
        paragraph = document.add_paragraph(str(item), style=style)
        overrides = element.get("style_overrides")
        if isinstance(overrides, dict):
            _apply_docx_v2_paragraph_spec(
                paragraph,
                overrides,
                {
                    "left": WD_ALIGN_PARAGRAPH.LEFT,
                    "center": WD_ALIGN_PARAGRAPH.CENTER,
                    "right": WD_ALIGN_PARAGRAPH.RIGHT,
                    "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
                },
            )
            _apply_docx_v2_run_spec(paragraph.runs, overrides)
    return len(items)


def _resolve_docx_v2_image_path(raw_path: str) -> Path:
    candidate = Path(raw_path)
    if candidate.is_absolute():
        raise ValueError("validation_error: DOCX image path must be workspace-relative.")
    if candidate.suffix.casefold() not in DOCX_V2_IMAGE_SUFFIXES:
        raise ValueError("validation_error: DOCX image path must use .png, .jpg or .jpeg.")
    parts = [part for part in raw_path.replace("\\", "/").split("/") if part]
    if any(part == ".." for part in parts):
        raise ValueError("validation_error: DOCX image path must not contain parent traversal.")
    if any(part.casefold() in DOCX_V2_PROTECTED_SEGMENTS for part in parts):
        raise ValueError("protected path is not allowed for DOCX images.")
    path = WORK_ROOT / candidate
    resolved = path.resolve(strict=True)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved != work_root and work_root not in resolved.parents:
        raise ValueError("DOCX image path must stay under /work.")
    if path.is_symlink() or resolved.is_symlink():
        raise ValueError("DOCX image symlinks are not supported.")
    if not resolved.is_file():
        raise ValueError("validation_error: DOCX image path must be a file.")
    if resolved.stat().st_size > DOCX_V2_MAX_IMAGE_BYTES:
        raise ValueError(f"DOCX image file exceeds limit of {DOCX_V2_MAX_IMAGE_BYTES} bytes.")
    try:
        from PIL import Image

        with Image.open(resolved) as image:
            if image.format not in {"PNG", "JPEG"}:
                raise ValueError("validation_error: DOCX image must be PNG or JPEG.")
            image.verify()
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("validation_error: DOCX image is not a readable PNG or JPEG.") from exc
    return resolved


def _docx_v2_limits() -> dict[str, int]:
    return {
        "max_json_input_bytes": DOCX_V2_MAX_JSON_INPUT_BYTES,
        "max_content_elements": DOCX_V2_MAX_CONTENT_ELEMENTS,
        "max_paragraphs": DOCX_V2_MAX_PARAGRAPHS,
        "max_tables": DOCX_V2_MAX_TABLES,
        "max_table_cells": DOCX_V2_MAX_TABLE_CELLS,
        "max_images": DOCX_V2_MAX_IMAGES,
        "max_image_bytes": DOCX_V2_MAX_IMAGE_BYTES,
    }


def _create_xlsx(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    if _is_xlsx_v2_content(content):
        return _create_xlsx_v2(output, title, content)
    return _create_xlsx_legacy(output, title, body, content)


def _is_xlsx_v2_content(content: dict[str, object]) -> bool:
    return any(key in content for key in ("schema_version", "workbook", "sheets"))


def _create_xlsx_legacy(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    from openpyxl import Workbook

    sheet_name = str(content.get("sheet_name") or "PLwC Smoke")[:31]
    rows = content.get("rows")
    row_values = rows if isinstance(rows, list) and rows else [["Title", "Value"], [title, 1], [body, 2], ["Total", 3]]
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = sheet_name or "PLwC Smoke"
    for row_index, row in enumerate(row_values, start=1):
        if not isinstance(row, list):
            raise ValueError("XLSX rows must be arrays.")
        for col_index, value in enumerate(row, start=1):
            _assign_xlsx_v2_value(sheet.cell(row=row_index, column=col_index), value)
    workbook.save(output)
    return {
        "builder_version": "xlsx_legacy_compat",
        "document_format": "xlsx",
        "sheet_count": 1,
        "sheet_names": [sheet.title],
        "total_row_count": len(row_values),
        "total_cell_count": sum(len(row) for row in row_values if isinstance(row, list)),
        "formula_count": 0,
    }


def _create_xlsx_v2(output: Path, title: str, content: dict[str, object]) -> dict[str, object]:
    from openpyxl import Workbook

    stats = _validate_xlsx_v2_worker_content(content)
    workbook = Workbook()
    metadata = content.get("workbook") if isinstance(content.get("workbook"), dict) else {}
    workbook.properties.title = str(metadata.get("title") or title)
    if isinstance(metadata.get("author"), str):
        workbook.properties.creator = str(metadata["author"])

    sheets = content["sheets"]
    for sheet_index, sheet_spec in enumerate(sheets):
        assert isinstance(sheet_spec, dict)
        sheet = workbook.active if sheet_index == 0 else workbook.create_sheet()
        sheet.title = str(sheet_spec["name"])
        rows = sheet_spec["rows"]
        assert isinstance(rows, list)
        for row_index, row in enumerate(rows, start=1):
            assert isinstance(row, list)
            for col_index, cell_spec in enumerate(row, start=1):
                cell = sheet.cell(row=row_index, column=col_index)
                if isinstance(cell_spec, dict):
                    if "formula" in cell_spec:
                        cell.value = "=" + _normalize_xlsx_v2_formula(str(cell_spec["formula"]))
                    else:
                        _assign_xlsx_v2_value(cell, cell_spec.get("value"))
                    _apply_xlsx_v2_cell_style(cell, cell_spec)
                else:
                    _assign_xlsx_v2_value(cell, cell_spec)

        _apply_xlsx_v2_sheet_structure(sheet, sheet_spec, max(len(row) for row in rows if isinstance(row, list)), len(rows))

    workbook.save(output)
    return {
        "builder_version": "xlsx_creation_v2",
        "document_format": "xlsx",
        "sheet_count": len(sheets),
        "sheet_names": [str(sheet["name"]) for sheet in sheets if isinstance(sheet, dict)],
        "total_row_count": stats["total_rows"],
        "total_cell_count": stats["total_cells"],
        "formula_count": stats["formula_count"],
        "merged_range_count": stats["merged_range_count"],
        "formula_policy": "Formulas are written only from explicit formula fields and are not executed by PLwC; value strings beginning with = are stored as literal text.",
        "limits": _xlsx_v2_limits(),
    }


def _assign_xlsx_v2_value(cell: object, value: object) -> None:
    cell.value = value
    if isinstance(value, str) and value.startswith("="):
        cell.data_type = "s"


def _apply_xlsx_v2_sheet_structure(sheet: object, sheet_spec: dict[str, object], max_cols: int, row_count: int) -> None:
    from openpyxl.utils import get_column_letter

    freeze_panes = sheet_spec.get("freeze_panes")
    if isinstance(freeze_panes, str):
        sheet.freeze_panes = freeze_panes.upper()
    column_widths = sheet_spec.get("column_widths")
    if isinstance(column_widths, dict):
        for column, width in column_widths.items():
            sheet.column_dimensions[str(column).upper()].width = float(width)
    row_heights = sheet_spec.get("row_heights")
    if isinstance(row_heights, dict):
        for row_number, height in row_heights.items():
            sheet.row_dimensions[int(row_number)].height = float(height)
    merge_cells = sheet_spec.get("merge_cells")
    if isinstance(merge_cells, list):
        for cell_range in merge_cells:
            sheet.merge_cells(str(cell_range).upper())
    if sheet_spec.get("auto_filter") is True and row_count > 0 and max_cols > 0:
        sheet.auto_filter.ref = f"A1:{get_column_letter(max_cols)}{row_count}"


def _apply_xlsx_v2_cell_style(cell: object, spec: dict[str, object]) -> None:
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    font_kwargs: dict[str, object] = {}
    if "bold" in spec:
        font_kwargs["bold"] = bool(spec["bold"])
    if "italic" in spec:
        font_kwargs["italic"] = bool(spec["italic"])
    if "underline" in spec and bool(spec["underline"]):
        font_kwargs["underline"] = "single"
    if "font_size" in spec:
        font_kwargs["size"] = float(spec["font_size"])
    if "font_color" in spec:
        font_kwargs["color"] = _normalize_xlsx_v2_color(str(spec["font_color"]))
    if font_kwargs:
        cell.font = Font(**font_kwargs)
    if "fill_color" in spec:
        cell.fill = PatternFill(fill_type="solid", fgColor=_normalize_xlsx_v2_color(str(spec["fill_color"])))
    alignment_kwargs: dict[str, object] = {}
    if "alignment" in spec:
        alignment_kwargs["horizontal"] = str(spec["alignment"]).casefold()
    if "vertical_alignment" in spec:
        vertical = str(spec["vertical_alignment"]).casefold()
        alignment_kwargs["vertical"] = "center" if vertical == "middle" else vertical
    if "wrap_text" in spec:
        alignment_kwargs["wrap_text"] = bool(spec["wrap_text"])
    if alignment_kwargs:
        cell.alignment = Alignment(**alignment_kwargs)
    if "number_format" in spec:
        cell.number_format = str(spec["number_format"])
    border = spec.get("border")
    if isinstance(border, str) and border.casefold() != "none":
        side = Side(style=border.casefold(), color="000000")
        cell.border = Border(left=side, right=side, top=side, bottom=side)


def _validate_xlsx_v2_worker_content(content: dict[str, object]) -> dict[str, int]:
    unknown_top = sorted(str(key) for key in content if str(key) not in XLSX_V2_TOP_LEVEL_FIELDS)
    if unknown_top:
        raise ValueError(f"validation_error: create_xlsx unsupported fields: {', '.join(unknown_top)}.")
    schema_version = content.get("schema_version", XLSX_V2_SCHEMA_VERSION)
    if not isinstance(schema_version, str) or schema_version not in XLSX_V2_SUPPORTED_SCHEMA_VERSIONS:
        raise ValueError(f"validation_error: create_xlsx schema_version must be one of: {', '.join(sorted(XLSX_V2_SUPPORTED_SCHEMA_VERSIONS))}.")
    workbook = content.get("workbook")
    if workbook is not None:
        if not isinstance(workbook, dict):
            raise ValueError("validation_error: create_xlsx workbook must be an object.")
        unknown_workbook = sorted(str(key) for key in workbook if str(key) not in XLSX_V2_WORKBOOK_FIELDS)
        if unknown_workbook:
            raise ValueError(f"validation_error: create_xlsx workbook has unsupported fields: {', '.join(unknown_workbook)}.")
        for field_name, value in workbook.items():
            if value is not None and not isinstance(value, str):
                raise ValueError(f"validation_error: create_xlsx workbook.{field_name} must be a string.")
    sheets = content.get("sheets")
    if not isinstance(sheets, list) or not sheets:
        raise ValueError("validation_error: create_xlsx sheets must be a non-empty list.")
    if len(sheets) > XLSX_V2_MAX_SHEETS:
        raise ValueError(f"create_xlsx sheet count exceeds {XLSX_V2_MAX_SHEETS}.")

    seen_names: set[str] = set()
    total_rows = 0
    total_cells = 0
    formula_count = 0
    merged_range_count = 0
    for sheet_index, sheet in enumerate(sheets, start=1):
        if not isinstance(sheet, dict):
            raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} must be an object.")
        sheet_stats = _validate_xlsx_v2_worker_sheet(sheet_index, sheet, seen_names)
        total_rows += sheet_stats["rows"]
        total_cells += sheet_stats["cells"]
        formula_count += sheet_stats["formulas"]
        merged_range_count += sheet_stats["merged_ranges"]
        if total_cells > XLSX_V2_MAX_TOTAL_CELLS:
            raise ValueError(f"create_xlsx total cell count exceeds {XLSX_V2_MAX_TOTAL_CELLS}.")
    return {
        "total_rows": total_rows,
        "total_cells": total_cells,
        "formula_count": formula_count,
        "merged_range_count": merged_range_count,
    }


def _validate_xlsx_v2_worker_sheet(sheet_index: int, sheet: dict[str, object], seen_names: set[str]) -> dict[str, int]:
    unknown = sorted(str(key) for key in sheet if str(key) not in XLSX_V2_SHEET_FIELDS)
    if unknown:
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} has unsupported fields: {', '.join(unknown)}.")
    name = sheet.get("name")
    if not _is_xlsx_v2_sheet_name(name):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.name is invalid.")
    normalized_name = str(name).casefold()
    if normalized_name in seen_names:
        raise ValueError(f"validation_error: create_xlsx duplicate sheet name: {name}.")
    seen_names.add(normalized_name)

    rows = sheet.get("rows")
    if not isinstance(rows, list) or not rows:
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.rows must be a non-empty list.")
    if len(rows) > XLSX_V2_MAX_ROWS_PER_SHEET:
        raise ValueError(f"create_xlsx sheet {sheet_index} row count exceeds {XLSX_V2_MAX_ROWS_PER_SHEET}.")

    cells = 0
    formulas = 0
    max_cols = 0
    for row_index, row in enumerate(rows, start=1):
        if not isinstance(row, list):
            raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} row {row_index} must be an array.")
        max_cols = max(max_cols, len(row))
        if max_cols > XLSX_V2_MAX_COLUMNS_PER_SHEET:
            raise ValueError(f"create_xlsx sheet {sheet_index} column count exceeds {XLSX_V2_MAX_COLUMNS_PER_SHEET}.")
        cells += len(row)
        if cells > XLSX_V2_MAX_CELLS_PER_SHEET:
            raise ValueError(f"create_xlsx sheet {sheet_index} cell count exceeds {XLSX_V2_MAX_CELLS_PER_SHEET}.")
        for col_index, cell in enumerate(row, start=1):
            if isinstance(cell, dict):
                formulas += _validate_xlsx_v2_worker_cell(sheet_index, row_index, col_index, cell)
            elif not _is_xlsx_v2_scalar(cell):
                raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index} has an unsupported value type.")

    _validate_xlsx_v2_worker_sheet_structure(sheet_index, sheet, row_count=len(rows), max_cols=max_cols)
    merged_ranges = len(sheet.get("merge_cells") or []) if isinstance(sheet.get("merge_cells"), list) else 0
    return {"rows": len(rows), "cells": cells, "formulas": formulas, "merged_ranges": merged_ranges}


def _validate_xlsx_v2_worker_cell(sheet_index: int, row_index: int, col_index: int, cell: dict[str, object]) -> int:
    unknown = sorted(str(key) for key in cell if str(key) not in XLSX_V2_CELL_FIELDS)
    if unknown:
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index} has unsupported fields: {', '.join(unknown)}.")
    has_value = "value" in cell
    has_formula = "formula" in cell
    if has_value and has_formula:
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index} must use either value or formula, not both.")
    if not has_value and not has_formula:
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index} requires value or formula.")
    if has_value and not _is_xlsx_v2_scalar(cell.get("value")):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.value has an unsupported type.")
    if has_formula:
        formula = cell.get("formula")
        if not isinstance(formula, str) or not formula.strip():
            raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.formula must be a non-empty string.")
        _validate_xlsx_v2_formula(formula)
    _validate_xlsx_v2_worker_cell_format(sheet_index, row_index, col_index, cell)
    return 1 if has_formula else 0


def _validate_xlsx_v2_worker_cell_format(sheet_index: int, row_index: int, col_index: int, cell: dict[str, object]) -> None:
    for field_name in ("bold", "italic", "underline", "wrap_text"):
        value = cell.get(field_name)
        if value is not None and not isinstance(value, bool):
            raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.{field_name} must be a boolean.")
    font_size = cell.get("font_size")
    if font_size is not None and (not _is_xlsx_v2_number(font_size) or not 1 <= float(font_size) <= 96):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.font_size must be a number between 1 and 96.")
    for field_name in ("font_color", "fill_color"):
        value = cell.get(field_name)
        if value is not None:
            _normalize_xlsx_v2_color(str(value))
    alignment = cell.get("alignment")
    if alignment is not None and (not isinstance(alignment, str) or alignment.casefold() not in XLSX_V2_ALIGNMENTS):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.alignment must be left, center or right.")
    vertical = cell.get("vertical_alignment")
    if vertical is not None and (not isinstance(vertical, str) or vertical.casefold() not in XLSX_V2_VERTICAL_ALIGNMENTS):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.vertical_alignment must be top, middle or bottom.")
    number_format = cell.get("number_format")
    if number_format is not None and (not isinstance(number_format, str) or not number_format.strip() or len(number_format) > 100):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.number_format must be a non-empty string up to 100 characters.")
    border = cell.get("border")
    if border is not None and (not isinstance(border, str) or border.casefold() not in XLSX_V2_BORDERS):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index} cell {row_index}:{col_index}.border must be none, thin, medium or thick.")


def _validate_xlsx_v2_worker_sheet_structure(sheet_index: int, sheet: dict[str, object], *, row_count: int, max_cols: int) -> None:
    freeze_panes = sheet.get("freeze_panes")
    if freeze_panes is not None and (not isinstance(freeze_panes, str) or not _is_xlsx_v2_cell_ref(freeze_panes)):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.freeze_panes must be an A1 cell reference.")
    auto_filter = sheet.get("auto_filter")
    if auto_filter is not None and not isinstance(auto_filter, bool):
        raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.auto_filter must be a boolean.")
    column_widths = sheet.get("column_widths")
    if column_widths is not None:
        if not isinstance(column_widths, dict):
            raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.column_widths must be an object.")
        for column, width in column_widths.items():
            if not isinstance(column, str) or not _is_xlsx_v2_column_ref(column):
                raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.column_widths keys must be Excel columns.")
            if not _is_xlsx_v2_number(width) or not 1 <= float(width) <= 100:
                raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.column_widths values must be numbers between 1 and 100.")
    row_heights = sheet.get("row_heights")
    if row_heights is not None:
        if not isinstance(row_heights, dict):
            raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.row_heights must be an object.")
        for row, height in row_heights.items():
            if not isinstance(row, str) or not row.isdigit() or int(row) < 1:
                raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.row_heights keys must be positive 1-based row numbers.")
            if not _is_xlsx_v2_number(height) or not 1 <= float(height) <= 300:
                raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.row_heights values must be numbers between 1 and 300.")
    merge_cells = sheet.get("merge_cells")
    if merge_cells is not None:
        if not isinstance(merge_cells, list):
            raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.merge_cells must be a list.")
        if len(merge_cells) > XLSX_V2_MAX_MERGED_RANGES:
            raise ValueError(f"create_xlsx merged range count exceeds {XLSX_V2_MAX_MERGED_RANGES}.")
        occupied: set[tuple[int, int]] = set()
        for item in merge_cells:
            if not isinstance(item, str) or not _is_xlsx_v2_range_ref(item):
                raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.merge_cells entries must be A1:B2 ranges.")
            min_col, min_row, max_col, max_row = _xlsx_v2_range_bounds(item)
            if max_row > row_count or max_col > max_cols:
                raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.merge_cells range exceeds used sheet bounds.")
            cells_in_range = (max_row - min_row + 1) * (max_col - min_col + 1)
            if cells_in_range <= 10000:
                current = {(row, col) for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
                if occupied & current:
                    raise ValueError(f"validation_error: create_xlsx sheet {sheet_index}.merge_cells ranges overlap.")
                occupied.update(current)


def _is_xlsx_v2_scalar(value: object) -> bool:
    return value is None or isinstance(value, (str, int, float, bool))


def _is_xlsx_v2_number(value: object) -> bool:
    return isinstance(value, (int, float)) and not isinstance(value, bool)


def _is_xlsx_v2_sheet_name(value: object) -> bool:
    return isinstance(value, str) and bool(value.strip()) and len(value) <= 31 and not re.search(r"[:\\/?*\[\]]", value)


def _is_xlsx_v2_column_ref(value: str) -> bool:
    match = re.fullmatch(r"[A-Za-z]{1,3}", value.strip())
    if not match:
        return False
    return 1 <= _xlsx_v2_column_index(value) <= 16384


def _is_xlsx_v2_cell_ref(value: str) -> bool:
    match = re.fullmatch(r"\$?([A-Za-z]{1,3})\$?([1-9][0-9]{0,6})", value.strip())
    if not match:
        return False
    col_index = _xlsx_v2_column_index(match.group(1))
    row_index = int(match.group(2))
    return 1 <= col_index <= 16384 and 1 <= row_index <= 1_048_576


def _is_xlsx_v2_range_ref(value: str) -> bool:
    parts = value.strip().split(":")
    if len(parts) != 2 or not all(_is_xlsx_v2_cell_ref(part) for part in parts):
        return False
    min_col, min_row, max_col, max_row = _xlsx_v2_range_bounds(value)
    return min_col <= max_col and min_row <= max_row


def _xlsx_v2_range_bounds(value: str) -> tuple[int, int, int, int]:
    start, end = value.strip().split(":", 1)
    start_col, start_row = _xlsx_v2_cell_ref_parts(start)
    end_col, end_row = _xlsx_v2_cell_ref_parts(end)
    return start_col, start_row, end_col, end_row


def _xlsx_v2_cell_ref_parts(value: str) -> tuple[int, int]:
    match = re.fullmatch(r"\$?([A-Za-z]{1,3})\$?([1-9][0-9]{0,6})", value.strip())
    if not match:
        raise ValueError("validation_error: invalid XLSX cell reference.")
    return _xlsx_v2_column_index(match.group(1)), int(match.group(2))


def _xlsx_v2_column_index(value: str) -> int:
    index = 0
    for char in value.strip().upper():
        index = index * 26 + (ord(char) - ord("A") + 1)
    return index


def _normalize_xlsx_v2_color(value: str) -> str:
    text = value.strip().removeprefix("#").upper()
    if not re.fullmatch(r"[0-9A-F]{6}|[0-9A-F]{8}", text):
        raise ValueError("validation_error: XLSX colors must be #RRGGBB or RRGGBB hex strings.")
    return text


def _normalize_xlsx_v2_formula(value: str) -> str:
    return value.strip().removeprefix("=")


def _validate_xlsx_v2_formula(value: str) -> None:
    formula = _normalize_xlsx_v2_formula(value)
    upper = formula.upper()
    if not formula:
        raise ValueError("validation_error: create_xlsx formula must not be empty.")
    if "[" in formula or "]" in formula or "://" in formula or "HTTP://" in upper or "HTTPS://" in upper:
        raise ValueError("validation_error: create_xlsx formulas must not contain external references or URLs.")
    if any(token in upper for token in XLSX_V2_UNSAFE_FORMULA_TOKENS):
        raise ValueError("validation_error: create_xlsx formula uses an unsafe external or macro-like function.")


def _xlsx_v2_limits() -> dict[str, int]:
    return {
        "max_json_input_bytes": XLSX_V2_MAX_JSON_INPUT_BYTES,
        "max_sheets": XLSX_V2_MAX_SHEETS,
        "max_rows_per_sheet": XLSX_V2_MAX_ROWS_PER_SHEET,
        "max_columns_per_sheet": XLSX_V2_MAX_COLUMNS_PER_SHEET,
        "max_cells_per_sheet": XLSX_V2_MAX_CELLS_PER_SHEET,
        "max_total_cells": XLSX_V2_MAX_TOTAL_CELLS,
        "max_merged_ranges": XLSX_V2_MAX_MERGED_RANGES,
    }


def _create_pptx(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    if _is_pptx_v2_content(content):
        return _create_pptx_v2(output, title, content)
    return _create_pptx_legacy(output, title, body, content)


def _is_pptx_v2_content(content: dict[str, object]) -> bool:
    if any(key in content for key in ("schema_version", "presentation")):
        return True
    slides = content.get("slides")
    if isinstance(slides, list):
        for slide in slides:
            if isinstance(slide, dict) and any(
                key in slide for key in ("layout", "subtitle", "body", "content", "notes", "image")
            ):
                return True
    return False


def _create_pptx_legacy(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    from pptx import Presentation

    deck_title = str(content.get("title") or title)
    slides = content.get("slides")
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = deck_title
    title_slide.placeholders[1].text = body
    slide_count = 1
    bullet_total = 0
    if isinstance(slides, list) and slides:
        for slide_item in slides:
            if not isinstance(slide_item, dict):
                raise ValueError("PPTX slides must be objects.")
            content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            content_slide.shapes.title.text = str(slide_item.get("title") or "Slide")
            bullets = _string_list(slide_item.get("bullets"), [])
            content_slide.placeholders[1].text = "\n".join(bullets)
            slide_count += 1
            bullet_total += len(bullets)
    else:
        content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        content_slide.shapes.title.text = "Document Worker MVP"
        content_slide.placeholders[1].text = "Creation smoke test only. Conversion is not implemented."
        slide_count += 1
    presentation.save(output)
    return {
        "builder_version": "pptx_legacy_compat",
        "document_format": "pptx",
        "slide_count": slide_count,
        "bullet_item_count": bullet_total,
    }


def _create_pptx_v2(output: Path, title: str, content: dict[str, object]) -> dict[str, object]:
    from pptx import Presentation
    from pptx.util import Mm

    stats = _validate_pptx_v2_worker_content(content)
    presentation_meta = content.get("presentation") if isinstance(content.get("presentation"), dict) else {}
    deck_title = str(presentation_meta.get("title") or title or "PLwC Presentation")
    author = presentation_meta.get("author")
    slide_size_spec = presentation_meta.get("slide_size", "16:9")

    presentation = Presentation()
    presentation.core_properties.title = deck_title
    if isinstance(author, str) and author.strip():
        presentation.core_properties.author = author

    slide_size_name, width_mm, height_mm = _resolve_pptx_v2_slide_size(slide_size_spec)
    presentation.slide_width = Mm(width_mm)
    presentation.slide_height = Mm(height_mm)
    blank_layout = presentation.slide_layouts[6]  # blank layout is index 6 in the python-pptx default

    slides_spec = content["slides"]
    notes_slide_count = 0
    for slide_spec in slides_spec:
        assert isinstance(slide_spec, dict)
        slide = presentation.slides.add_slide(blank_layout)
        layout = slide_spec["layout"]
        if layout == "title":
            _render_pptx_v2_title_slide(slide, slide_spec, width_mm, height_mm)
        elif layout == "section_header":
            _render_pptx_v2_section_header_slide(slide, slide_spec, width_mm, height_mm)
        elif layout == "content":
            _render_pptx_v2_content_slide(slide, slide_spec, width_mm, height_mm)
        elif layout == "image":
            _render_pptx_v2_image_slide(slide, slide_spec, width_mm, height_mm)
        elif layout == "blank":
            _render_pptx_v2_blank_slide(slide, slide_spec, width_mm, height_mm)
        notes = slide_spec.get("notes")
        if isinstance(notes, str) and notes:
            _set_pptx_v2_notes(slide, notes)
            notes_slide_count += 1

    presentation.save(output)
    return {
        "builder_version": "pptx_creation_v2",
        "document_format": "pptx",
        "slide_count": stats["slide_count"],
        "content_element_count": stats["content_element_count"],
        "bullet_item_count": stats["bullet_item_count"],
        "table_count": stats["table_count"],
        "table_cell_count": stats["table_cell_count"],
        "image_count": stats["image_count"],
        "text_box_count": stats["text_box_count"],
        "paragraph_count": stats["paragraph_count"],
        "notes_slide_count": notes_slide_count,
        "slide_size": slide_size_name,
        "slide_width_mm": width_mm,
        "slide_height_mm": height_mm,
        "limits": _pptx_v2_limits(),
    }


def _resolve_pptx_v2_slide_size(slide_size: object) -> tuple[str, float, float]:
    if isinstance(slide_size, str):
        if slide_size not in PPTX_V2_SLIDE_SIZE_PRESETS_MM:
            raise ValueError(
                "validation_error: PPTX slide_size must be 16:9, 4:3, A4_portrait or a custom object."
            )
        width, height = PPTX_V2_SLIDE_SIZE_PRESETS_MM[slide_size]
        return slide_size, float(width), float(height)
    if isinstance(slide_size, dict):
        width = slide_size.get("width_mm")
        height = slide_size.get("height_mm")
        if (
            not isinstance(width, (int, float))
            or isinstance(width, bool)
            or not isinstance(height, (int, float))
            or isinstance(height, bool)
        ):
            raise ValueError(
                "validation_error: PPTX custom slide_size must define numeric width_mm and height_mm."
            )
        name = slide_size.get("name") or "custom"
        return str(name), float(width), float(height)
    raise ValueError("validation_error: PPTX slide_size must be a string preset or a custom object.")


# --- Layout renderers -------------------------------------------------------


def _render_pptx_v2_title_slide(slide: object, slide_spec: dict[str, object], width_mm: float, height_mm: float) -> None:
    title_box = _add_text_frame(slide, left_mm=20, top_mm=height_mm * 0.35, width_mm=width_mm - 40, height_mm=30)
    _set_text_frame_text(title_box, str(slide_spec.get("title") or ""), bold=True, font_size=44, alignment="center")
    subtitle = slide_spec.get("subtitle")
    if isinstance(subtitle, str) and subtitle.strip():
        sub_box = _add_text_frame(
            slide, left_mm=20, top_mm=height_mm * 0.35 + 32, width_mm=width_mm - 40, height_mm=20
        )
        _set_text_frame_text(sub_box, subtitle, font_size=24, alignment="center")


def _render_pptx_v2_section_header_slide(
    slide: object, slide_spec: dict[str, object], width_mm: float, height_mm: float
) -> None:
    title_box = _add_text_frame(slide, left_mm=20, top_mm=height_mm * 0.4, width_mm=width_mm - 40, height_mm=30)
    _set_text_frame_text(title_box, str(slide_spec.get("title") or ""), bold=True, font_size=40, alignment="left")
    body = slide_spec.get("body")
    if isinstance(body, str) and body.strip():
        body_box = _add_text_frame(
            slide, left_mm=20, top_mm=height_mm * 0.4 + 32, width_mm=width_mm - 40, height_mm=30
        )
        _set_text_frame_text(body_box, body, font_size=18, alignment="left")


def _render_pptx_v2_content_slide(
    slide: object, slide_spec: dict[str, object], width_mm: float, height_mm: float
) -> None:
    title_box = _add_text_frame(slide, left_mm=15, top_mm=12, width_mm=width_mm - 30, height_mm=18)
    _set_text_frame_text(title_box, str(slide_spec.get("title") or ""), bold=True, font_size=32, alignment="left")
    content_elements = slide_spec.get("content")
    if not isinstance(content_elements, list):
        return
    cursor_top = 36.0
    content_left = 15.0
    content_width = width_mm - 30
    content_bottom_limit = height_mm - 12
    for element in content_elements:
        if not isinstance(element, dict):
            continue
        element_type = element.get("type")
        remaining_height = max(content_bottom_limit - cursor_top, 10)
        if element_type == "bullets":
            block = _add_text_frame(
                slide, left_mm=content_left, top_mm=cursor_top, width_mm=content_width, height_mm=remaining_height
            )
            _render_pptx_v2_bullets(block, element.get("items") or [])
            cursor_top += min(remaining_height, max(20, len(element.get("items") or []) * 8))
        elif element_type == "paragraph":
            block = _add_text_frame(
                slide, left_mm=content_left, top_mm=cursor_top, width_mm=content_width, height_mm=20
            )
            _set_text_frame_text(
                block,
                str(element.get("text") or ""),
                bold=bool(element.get("bold")),
                italic=bool(element.get("italic")),
                font_size=element.get("font_size"),
                color=element.get("color"),
                alignment=element.get("alignment"),
            )
            cursor_top += 20
        elif element_type == "table":
            rows = element.get("rows") or []
            header_row = bool(element.get("header_row"))
            row_count = len(rows)
            col_count = len(rows[0]) if rows and isinstance(rows[0], list) else 0
            if row_count == 0 or col_count == 0:
                continue
            table_height = max(10.0, min(remaining_height, row_count * 9.0))
            _render_pptx_v2_table(
                slide,
                rows,
                header_row=header_row,
                left_mm=content_left,
                top_mm=cursor_top,
                width_mm=content_width,
                height_mm=table_height,
            )
            cursor_top += table_height + 4
        elif element_type == "image":
            _render_pptx_v2_inline_image(
                slide,
                element,
                left_mm=content_left,
                top_mm=cursor_top,
                width_mm=content_width,
                remaining_height_mm=remaining_height,
            )
            cursor_top += min(remaining_height, float(element.get("height_mm") or 40))
        elif element_type == "text_box":
            _render_pptx_v2_text_box(slide, element)


def _render_pptx_v2_image_slide(
    slide: object, slide_spec: dict[str, object], width_mm: float, height_mm: float
) -> None:
    title_text = slide_spec.get("title")
    title_height = 0.0
    if isinstance(title_text, str) and title_text.strip():
        title_box = _add_text_frame(slide, left_mm=15, top_mm=12, width_mm=width_mm - 30, height_mm=18)
        _set_text_frame_text(title_box, title_text, bold=True, font_size=28, alignment="left")
        title_height = 24.0
    image_spec = slide_spec.get("image")
    if not isinstance(image_spec, dict):
        return
    _render_pptx_v2_inline_image(
        slide,
        image_spec,
        left_mm=15,
        top_mm=12 + title_height,
        width_mm=width_mm - 30,
        remaining_height_mm=height_mm - (12 + title_height) - 12,
    )


def _render_pptx_v2_blank_slide(
    slide: object, slide_spec: dict[str, object], width_mm: float, height_mm: float
) -> None:
    content_elements = slide_spec.get("content")
    if not isinstance(content_elements, list):
        return
    for element in content_elements:
        if not isinstance(element, dict):
            continue
        if element.get("type") == "text_box":
            _render_pptx_v2_text_box(slide, element)


# --- Element renderers ------------------------------------------------------


def _render_pptx_v2_bullets(text_frame: object, items: list[object]) -> None:
    text_frame.word_wrap = True
    first = True
    for item in items:
        if not isinstance(item, dict):
            continue
        text = str(item.get("text") or "")
        level = int(item.get("level") or 0)
        if first:
            paragraph = text_frame.paragraphs[0]
            for extra_run in list(paragraph.runs):
                extra_run.text = ""
            first = False
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.level = level
        if paragraph.runs:
            target_run = paragraph.runs[0]
            target_run.text = text
        else:
            target_run = paragraph.add_run()
            target_run.text = text
        _apply_pptx_v2_run_format(
            target_run,
            bold=bool(item.get("bold")),
            italic=bool(item.get("italic")),
            font_size=item.get("font_size"),
            color=item.get("color"),
        )


def _render_pptx_v2_text_box(slide: object, element: dict[str, object]) -> None:
    left_mm = float(element.get("left_mm") or 0)
    top_mm = float(element.get("top_mm") or 0)
    width_mm = float(element.get("width_mm") or 100)
    height_mm = float(element.get("height_mm") or 30)
    text_frame = _add_text_frame(slide, left_mm=left_mm, top_mm=top_mm, width_mm=width_mm, height_mm=height_mm)
    _set_text_frame_text(
        text_frame,
        str(element.get("text") or ""),
        bold=bool(element.get("bold")),
        italic=bool(element.get("italic")),
        font_size=element.get("font_size"),
        color=element.get("color"),
        alignment=element.get("alignment"),
    )


def _render_pptx_v2_table(
    slide: object,
    rows: list[object],
    *,
    header_row: bool,
    left_mm: float,
    top_mm: float,
    width_mm: float,
    height_mm: float,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Mm

    row_count = len(rows)
    col_count = max(len(row) for row in rows if isinstance(row, list))
    table_shape = slide.shapes.add_table(row_count, col_count, Mm(left_mm), Mm(top_mm), Mm(width_mm), Mm(height_mm))
    table = table_shape.table
    for row_index, row in enumerate(rows):
        if not isinstance(row, list):
            continue
        for col_index, cell_value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = "" if cell_value is None else str(cell_value)
            if header_row and row_index == 0:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        try:
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        except Exception:
                            pass


def _render_pptx_v2_inline_image(
    slide: object,
    spec: dict[str, object],
    *,
    left_mm: float,
    top_mm: float,
    width_mm: float,
    remaining_height_mm: float,
) -> None:
    from pptx.util import Mm

    raw_path = spec.get("path")
    if not isinstance(raw_path, str):
        return
    resolved = _resolve_pptx_v2_image_path(raw_path)
    requested_width = spec.get("width_mm")
    requested_height = spec.get("height_mm")
    width = float(requested_width) if isinstance(requested_width, (int, float)) and not isinstance(requested_width, bool) else width_mm
    height = float(requested_height) if isinstance(requested_height, (int, float)) and not isinstance(requested_height, bool) else None
    if height is None:
        height = min(remaining_height_mm, width * 0.5625)
    align = str(spec.get("align") or "left").casefold()
    if align == "center":
        offset_left = left_mm + max(0.0, (width_mm - width) / 2)
    elif align == "right":
        offset_left = left_mm + max(0.0, width_mm - width)
    else:
        offset_left = left_mm
    slide.shapes.add_picture(str(resolved), Mm(offset_left), Mm(top_mm), Mm(width), Mm(height))


# --- Notes ------------------------------------------------------------------


def _set_pptx_v2_notes(slide: object, notes: str) -> None:
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


# --- Text frame helpers -----------------------------------------------------


def _add_text_frame(
    slide: object, *, left_mm: float, top_mm: float, width_mm: float, height_mm: float
) -> object:
    from pptx.util import Mm

    box = slide.shapes.add_textbox(Mm(left_mm), Mm(top_mm), Mm(width_mm), Mm(height_mm))
    return box.text_frame


def _set_text_frame_text(
    text_frame: object,
    text: str,
    *,
    bold: bool = False,
    italic: bool = False,
    font_size: object = None,
    color: object = None,
    alignment: object = None,
) -> None:
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    for extra_run in list(paragraph.runs)[1:]:
        extra_run.text = ""
    if paragraph.runs:
        run = paragraph.runs[0]
        run.text = text
    else:
        run = paragraph.add_run()
        run.text = text
    _apply_pptx_v2_paragraph_alignment(paragraph, alignment)
    _apply_pptx_v2_run_format(run, bold=bold, italic=italic, font_size=font_size, color=color)


def _apply_pptx_v2_paragraph_alignment(paragraph: object, alignment: object) -> None:
    from pptx.enum.text import PP_ALIGN

    if not isinstance(alignment, str):
        return
    value = alignment.casefold()
    mapping = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    if value in mapping:
        paragraph.alignment = mapping[value]


def _apply_pptx_v2_run_format(
    run: object,
    *,
    bold: bool = False,
    italic: bool = False,
    font_size: object = None,
    color: object = None,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True
    if isinstance(font_size, (int, float)) and not isinstance(font_size, bool):
        run.font.size = Pt(float(font_size))
    if isinstance(color, str) and color.strip():
        normalized = _normalize_pptx_v2_color(color)
        try:
            run.font.color.rgb = RGBColor(int(normalized[0:2], 16), int(normalized[2:4], 16), int(normalized[4:6], 16))
        except Exception:
            pass


def _normalize_pptx_v2_color(value: str) -> str:
    text = value.strip().lstrip("#").upper()
    if not re.fullmatch(r"[0-9A-F]{6}", text):
        raise ValueError("validation_error: PPTX colors must be #RRGGBB or RRGGBB hex strings.")
    return text


# --- Worker-side validation -------------------------------------------------


def _validate_pptx_v2_worker_content(content: dict[str, object]) -> dict[str, int]:
    unknown_top = sorted(str(key) for key in content if str(key) not in PPTX_V2_TOP_LEVEL_FIELDS)
    if unknown_top:
        raise ValueError(f"validation_error: create_pptx unsupported fields: {', '.join(unknown_top)}.")
    schema_version = content.get("schema_version", PPTX_V2_SCHEMA_VERSION)
    if not isinstance(schema_version, str) or schema_version not in PPTX_V2_SUPPORTED_SCHEMA_VERSIONS:
        raise ValueError(
            f"validation_error: create_pptx schema_version must be one of: "
            f"{', '.join(sorted(PPTX_V2_SUPPORTED_SCHEMA_VERSIONS))}."
        )
    presentation = content.get("presentation")
    if presentation is not None:
        if not isinstance(presentation, dict):
            raise ValueError("validation_error: create_pptx presentation must be an object.")
        unknown_presentation = sorted(str(key) for key in presentation if str(key) not in PPTX_V2_PRESENTATION_FIELDS)
        if unknown_presentation:
            raise ValueError(
                f"validation_error: create_pptx presentation has unsupported fields: "
                f"{', '.join(unknown_presentation)}."
            )
    slides = content.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("validation_error: create_pptx slides must be a non-empty list.")
    if len(slides) > PPTX_V2_MAX_SLIDES:
        raise ValueError(f"create_pptx slide count exceeds {PPTX_V2_MAX_SLIDES}.")
    stats = {
        "slide_count": len(slides),
        "content_element_count": 0,
        "bullet_item_count": 0,
        "table_count": 0,
        "table_cell_count": 0,
        "image_count": 0,
        "text_box_count": 0,
        "paragraph_count": 0,
    }
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            raise ValueError(f"validation_error: create_pptx slide {index} must be an object.")
        _validate_pptx_v2_worker_slide(index, slide, stats)
    if stats["image_count"] > PPTX_V2_MAX_IMAGES:
        raise ValueError(f"create_pptx image count exceeds {PPTX_V2_MAX_IMAGES}.")
    return stats


def _validate_pptx_v2_worker_slide(index: int, slide: dict[str, object], stats: dict[str, int]) -> None:
    unknown = sorted(str(key) for key in slide if str(key) not in PPTX_V2_SLIDE_FIELDS)
    if unknown:
        raise ValueError(f"validation_error: create_pptx slide {index} has unsupported fields: {', '.join(unknown)}.")
    layout = slide.get("layout")
    if not isinstance(layout, str) or layout not in PPTX_V2_LAYOUTS:
        raise ValueError(
            f"validation_error: create_pptx slide {index}.layout must be one of: {', '.join(sorted(PPTX_V2_LAYOUTS))}."
        )
    notes = slide.get("notes")
    if notes is not None and not isinstance(notes, str):
        raise ValueError(f"validation_error: create_pptx slide {index}.notes must be a string.")
    if layout == "image":
        image_spec = slide.get("image")
        if not isinstance(image_spec, dict):
            raise ValueError(f"validation_error: create_pptx slide {index} (image layout) requires image object.")
        _validate_pptx_v2_worker_image_spec(index, image_spec)
        stats["image_count"] += 1
    elif layout in {"content", "blank"}:
        elements = slide.get("content")
        if elements is not None:
            if not isinstance(elements, list) or not elements:
                raise ValueError(
                    f"validation_error: create_pptx slide {index}.content must be a non-empty list when present."
                )
            if len(elements) > PPTX_V2_MAX_CONTENT_ELEMENTS_PER_SLIDE:
                raise ValueError(
                    f"create_pptx slide {index} content element count exceeds {PPTX_V2_MAX_CONTENT_ELEMENTS_PER_SLIDE}."
                )
            slide_bullets = 0
            slide_tables = 0
            slide_table_cells = 0
            for element_index, element in enumerate(elements, start=1):
                if not isinstance(element, dict):
                    raise ValueError(
                        f"validation_error: create_pptx slide {index} content element {element_index} must be an object."
                    )
                element_type = element.get("type")
                if not isinstance(element_type, str) or element_type not in PPTX_V2_CONTENT_TYPES:
                    raise ValueError(
                        f"validation_error: create_pptx slide {index} content element {element_index}.type must be one of: "
                        f"{', '.join(sorted(PPTX_V2_CONTENT_TYPES))}."
                    )
                if element_type == "bullets":
                    item_count = _validate_pptx_v2_worker_bullets(index, element_index, element)
                    slide_bullets += item_count
                    if slide_bullets > PPTX_V2_MAX_BULLET_ITEMS_PER_SLIDE:
                        raise ValueError(
                            f"create_pptx slide {index} bullet item count exceeds {PPTX_V2_MAX_BULLET_ITEMS_PER_SLIDE}."
                        )
                    stats["bullet_item_count"] += item_count
                elif element_type == "paragraph":
                    stats["paragraph_count"] += 1
                elif element_type == "table":
                    slide_tables += 1
                    if slide_tables > PPTX_V2_MAX_TABLES_PER_SLIDE:
                        raise ValueError(
                            f"create_pptx slide {index} table count exceeds {PPTX_V2_MAX_TABLES_PER_SLIDE}."
                        )
                    cell_count = _validate_pptx_v2_worker_table(index, element_index, element)
                    slide_table_cells += cell_count
                    if slide_table_cells > PPTX_V2_MAX_TABLE_CELLS_PER_SLIDE:
                        raise ValueError(
                            f"create_pptx slide {index} table cell count exceeds {PPTX_V2_MAX_TABLE_CELLS_PER_SLIDE}."
                        )
                    stats["table_count"] += 1
                    stats["table_cell_count"] += cell_count
                elif element_type == "image":
                    _validate_pptx_v2_worker_image_spec(index, element)
                    stats["image_count"] += 1
                elif element_type == "text_box":
                    stats["text_box_count"] += 1
                stats["content_element_count"] += 1


def _validate_pptx_v2_worker_bullets(slide_index: int, element_index: int, element: dict[str, object]) -> int:
    items = element.get("items")
    if not isinstance(items, list) or not items:
        raise ValueError(
            f"validation_error: create_pptx slide {slide_index} bullets element {element_index}.items must be a non-empty list."
        )
    for item_index, item in enumerate(items, start=1):
        if not isinstance(item, dict):
            raise ValueError(
                f"validation_error: create_pptx slide {slide_index} bullets element {element_index} item {item_index} must be an object."
            )
        text = item.get("text")
        if not isinstance(text, str) or not text.strip():
            raise ValueError(
                f"validation_error: create_pptx slide {slide_index} bullets element {element_index} item {item_index}.text must be a non-empty string."
            )
        level = item.get("level", 0)
        if isinstance(level, bool) or not isinstance(level, int) or level not in PPTX_V2_BULLET_LEVELS:
            raise ValueError(
                f"validation_error: create_pptx slide {slide_index} bullets element {element_index} item {item_index}.level "
                f"must be one of {sorted(PPTX_V2_BULLET_LEVELS)}."
            )
    return len(items)


def _validate_pptx_v2_worker_table(slide_index: int, element_index: int, element: dict[str, object]) -> int:
    rows = element.get("rows")
    if not isinstance(rows, list) or not rows:
        raise ValueError(
            f"validation_error: create_pptx slide {slide_index} table element {element_index}.rows must be a non-empty list."
        )
    column_count = None
    cell_count = 0
    for row_index, row in enumerate(rows, start=1):
        if not isinstance(row, list) or not all(
            cell is None or isinstance(cell, (str, int, float, bool)) for cell in row
        ):
            raise ValueError(
                f"validation_error: create_pptx slide {slide_index} table element {element_index} row {row_index} must be a list of scalar cells."
            )
        if column_count is None:
            column_count = len(row)
        elif len(row) != column_count:
            raise ValueError(
                f"validation_error: create_pptx slide {slide_index} table element {element_index} rows must share the same column count."
            )
        cell_count += len(row)
    return cell_count


def _validate_pptx_v2_worker_image_spec(slide_index: int, spec: dict[str, object]) -> None:
    path_value = spec.get("path")
    if not isinstance(path_value, str) or not path_value.strip():
        raise ValueError(f"validation_error: create_pptx slide {slide_index} image requires a workspace-relative path.")
    _resolve_pptx_v2_image_path(path_value)


def _resolve_pptx_v2_image_path(raw_path: str) -> Path:
    text = raw_path.strip()
    if re.match(r"^[A-Za-z][A-Za-z0-9+.\-]+:", text):
        raise ValueError("validation_error: PPTX image path must not be a URL.")
    if text.startswith("\\\\") or text.startswith("//"):
        raise ValueError("validation_error: PPTX image path must not be a UNC network path.")
    candidate = Path(text)
    if candidate.is_absolute():
        raise ValueError("validation_error: PPTX image path must be workspace-relative.")
    parts = [part for part in text.replace("\\", "/").split("/") if part]
    if any(part == ".." for part in parts):
        raise ValueError("validation_error: PPTX image path must not contain parent traversal.")
    if any(part.casefold() in PPTX_V2_PROTECTED_SEGMENTS for part in parts):
        raise ValueError("validation_error: PPTX image path must not target profile or governance directories.")
    if candidate.suffix.casefold() not in PPTX_V2_IMAGE_SUFFIXES:
        raise ValueError("validation_error: PPTX image path must use .png, .jpg or .jpeg.")
    path = WORK_ROOT / candidate
    resolved = path.resolve(strict=True)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved != work_root and work_root not in resolved.parents:
        raise ValueError("validation_error: PPTX image path must stay under /work.")
    if path.is_symlink() or resolved.is_symlink():
        raise ValueError("validation_error: PPTX image symlinks are not supported.")
    if not resolved.is_file():
        raise ValueError("validation_error: PPTX image path must be a file.")
    if resolved.stat().st_size > PPTX_V2_MAX_IMAGE_BYTES:
        raise ValueError(f"PPTX image file exceeds limit of {PPTX_V2_MAX_IMAGE_BYTES} bytes.")
    try:
        from PIL import Image

        with Image.open(resolved) as image:
            if image.format not in {"PNG", "JPEG"}:
                raise ValueError("validation_error: PPTX image must be PNG or JPEG.")
            image.verify()
    except ValueError:
        raise
    except Exception as exc:  # pragma: no cover - delegated to PIL
        raise ValueError("validation_error: PPTX image is not a readable PNG or JPEG.") from exc
    return resolved


def _pptx_v2_limits() -> dict[str, int]:
    return {
        "max_json_input_bytes": PPTX_V2_MAX_JSON_INPUT_BYTES,
        "max_slides": PPTX_V2_MAX_SLIDES,
        "max_content_elements_per_slide": PPTX_V2_MAX_CONTENT_ELEMENTS_PER_SLIDE,
        "max_bullet_items_per_slide": PPTX_V2_MAX_BULLET_ITEMS_PER_SLIDE,
        "max_tables_per_slide": PPTX_V2_MAX_TABLES_PER_SLIDE,
        "max_table_cells_per_slide": PPTX_V2_MAX_TABLE_CELLS_PER_SLIDE,
        "max_images": PPTX_V2_MAX_IMAGES,
        "max_image_bytes": PPTX_V2_MAX_IMAGE_BYTES,
    }


def _create_pdf(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    if _is_pdf_v2_content(content):
        return _create_pdf_v2(output, title, content)
    return _create_pdf_legacy(output, title, body, content)


def _is_pdf_v2_content(content: dict[str, object]) -> bool:
    return any(key in content for key in ("schema_version", "document", "page", "styles", "content"))


def _create_pdf_legacy(output: Path, title: str, body: str, content: dict[str, object]) -> dict[str, object]:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    document_title = str(content.get("title") or title)
    lines = _string_list(content.get("lines"), [body, "Page 1"])
    pdf = canvas.Canvas(str(output), pagesize=letter)
    pdf.setTitle(document_title)
    pdf.drawString(72, 720, document_title)
    y = 700
    for line in lines:
        pdf.drawString(72, y, line)
        y -= 20
    pdf.showPage()
    pdf.save()
    return {
        "builder_version": "pdf_legacy_compat",
        "document_format": "pdf",
        "content_element_count": len(lines) + 1,
        "paragraph_count": len(lines) + 1,
        "table_count": 0,
        "image_count": 0,
    }


# ---------------------------------------------------------------------------
# PDF Creation V2 layout builder
# ---------------------------------------------------------------------------


def _create_pdf_v2(output: Path, title: str, content: dict[str, object]) -> dict[str, object]:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT, TA_RIGHT
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        Image,
        ListFlowable,
        ListItem,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    stats = _validate_pdf_v2_worker_content(content)
    metadata = content.get("document") if isinstance(content.get("document"), dict) else {}
    document_title = str(metadata.get("title") or title or "PLwC PDF")
    document_author = str(metadata.get("author") or "")
    document_language = str(metadata.get("language") or "")

    page_settings = _pdf_v2_page_settings(content.get("page"))
    width_mm, height_mm = page_settings["width_mm"], page_settings["height_mm"]
    page_size = (width_mm * mm, height_mm * mm)

    margins = page_settings["margins_mm"]
    doc = SimpleDocTemplate(
        str(output),
        pagesize=page_size,
        leftMargin=margins["left"] * mm,
        rightMargin=margins["right"] * mm,
        topMargin=margins["top"] * mm,
        bottomMargin=margins["bottom"] * mm,
        title=document_title,
        author=document_author,
    )

    style_specs = content.get("styles") if isinstance(content.get("styles"), dict) else {}
    base_styles = getSampleStyleSheet()
    pdf_styles = _build_pdf_v2_styles(base_styles, style_specs)

    flowables: list[object] = []
    align_map = {"left": TA_LEFT, "center": TA_CENTER, "right": TA_RIGHT, "justify": TA_JUSTIFY}

    for element in content["content"]:
        assert isinstance(element, dict)
        element_type = element["type"]
        if element_type in PDF_V2_TEXT_TYPES:
            paragraph_text = _pdf_v2_paragraph_text(element)
            style = pdf_styles.get(element_type, pdf_styles["body"])
            overrides = element.get("style_overrides")
            if isinstance(overrides, dict):
                style = _apply_pdf_v2_style_overrides(style, overrides, align_map)
            if element_type == "heading1" and pdf_styles["heading1_breaks_before"] and flowables:
                flowables.append(PageBreak())
            if paragraph_text:
                flowables.append(Paragraph(paragraph_text, style))
            else:
                # Empty paragraph: emit a small spacer to preserve structure
                flowables.append(Spacer(1, 4 * mm))
        elif element_type in {"bullet_list", "numbered_list"}:
            items_style = pdf_styles["body"]
            overrides = element.get("style_overrides")
            if isinstance(overrides, dict):
                items_style = _apply_pdf_v2_style_overrides(items_style, overrides, align_map)
            bullet_type = "bullet" if element_type == "bullet_list" else "1"
            list_items = [
                ListItem(Paragraph(_pdf_v2_escape(str(item)), items_style))
                for item in element["items"]
            ]
            flowables.append(
                ListFlowable(
                    list_items,
                    bulletType=bullet_type,
                    bulletFontName=items_style.fontName,
                    leftIndent=12,
                )
            )
        elif element_type == "page_break":
            flowables.append(PageBreak())
        elif element_type == "image":
            resolved = _resolve_pdf_v2_image_path(str(element["path"]))
            image_width_mm = element.get("width_mm")
            image_height_mm = element.get("height_mm")
            image = Image(str(resolved))
            if isinstance(image_width_mm, (int, float)) and not isinstance(image_width_mm, bool):
                if isinstance(image_height_mm, (int, float)) and not isinstance(image_height_mm, bool):
                    image.drawWidth = float(image_width_mm) * mm
                    image.drawHeight = float(image_height_mm) * mm
                else:
                    ratio = float(image_width_mm) * mm / image.drawWidth
                    image.drawWidth = float(image_width_mm) * mm
                    image.drawHeight = image.drawHeight * ratio
            elif isinstance(image_height_mm, (int, float)) and not isinstance(image_height_mm, bool):
                ratio = float(image_height_mm) * mm / image.drawHeight
                image.drawHeight = float(image_height_mm) * mm
                image.drawWidth = image.drawWidth * ratio
            align = str(element.get("align") or "left").casefold()
            image.hAlign = {"left": "LEFT", "center": "CENTER", "right": "RIGHT"}.get(align, "LEFT")
            flowables.append(image)
        elif element_type == "table":
            rows = element["rows"]
            assert isinstance(rows, list)
            normalized_rows = [
                ["" if cell is None else str(cell) for cell in row] for row in rows
            ]
            table = Table(normalized_rows, repeatRows=1 if element.get("header_row") else 0)
            table_style_commands: list[tuple] = [
                ("FONTNAME", (0, 0), (-1, -1), pdf_styles["body"].fontName),
                ("FONTSIZE", (0, 0), (-1, -1), pdf_styles["body"].fontSize),
                ("BOX", (0, 0), (-1, -1), 0.5, colors.black),
                ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.grey),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
            if element.get("header_row"):
                table_style_commands.extend(
                    [
                        ("FONTNAME", (0, 0), (-1, 0), pdf_styles["heading2"].fontName),
                        ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.black),
                    ]
                )
            table.setStyle(TableStyle(table_style_commands))
            flowables.append(table)

    doc.build(flowables)
    _ = document_language  # retained for future metadata; reportlab has no setter.

    response = {
        "builder_version": "pdf_creation_v2",
        "document_format": "pdf",
        "content_element_count": stats["content_element_count"],
        "paragraph_count": stats["paragraph_count"],
        "table_count": stats["table_count"],
        "table_cell_count": stats["table_cell_count"],
        "image_count": stats["image_count"],
        "list_item_count": stats["list_item_count"],
        "page_size": page_settings["page_size"],
        "orientation": page_settings["orientation"],
        "page_width_mm": page_settings["width_mm"],
        "page_height_mm": page_settings["height_mm"],
        "limits": _pdf_v2_limits(),
    }
    return response


def _pdf_v2_paragraph_text(element: dict[str, object]) -> str:
    if isinstance(element.get("text"), str):
        return _pdf_v2_escape(element["text"])
    runs = element.get("runs")
    if isinstance(runs, list):
        parts: list[str] = []
        for run in runs:
            if not isinstance(run, dict):
                continue
            text = _pdf_v2_escape(str(run.get("text") or ""))
            if run.get("bold"):
                text = f"<b>{text}</b>"
            if run.get("italic"):
                text = f"<i>{text}</i>"
            if run.get("underline"):
                text = f"<u>{text}</u>"
            parts.append(text)
        return "".join(parts)
    return ""


def _pdf_v2_escape(value: str) -> str:
    return (
        value.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def _pdf_v2_page_settings(spec: object) -> dict[str, object]:
    settings = {
        "page_size": "A4",
        "orientation": "portrait",
        "width_mm": float(PDF_V2_PAGE_SIZES_MM["A4"][0]),
        "height_mm": float(PDF_V2_PAGE_SIZES_MM["A4"][1]),
        "margins_mm": {"top": 18.0, "bottom": 18.0, "left": 18.0, "right": 18.0},
    }
    if not isinstance(spec, dict):
        return settings
    size_value = spec.get("size", "A4")
    if isinstance(size_value, str):
        size_key = size_value.upper()
        if size_key not in PDF_V2_PAGE_SIZES_MM:
            raise ValueError("validation_error: create_pdf page.size must be A4, A5 or a custom object.")
        width, height = PDF_V2_PAGE_SIZES_MM[size_key]
        settings["page_size"] = size_key
        settings["width_mm"] = float(width)
        settings["height_mm"] = float(height)
    elif isinstance(size_value, dict):
        name = size_value.get("name") or "custom"
        width = size_value.get("width_mm")
        height = size_value.get("height_mm")
        if (
            not isinstance(width, (int, float))
            or isinstance(width, bool)
            or not isinstance(height, (int, float))
            or isinstance(height, bool)
        ):
            raise ValueError(
                "validation_error: create_pdf custom page.size requires numeric width_mm and height_mm."
            )
        settings["page_size"] = str(name)
        settings["width_mm"] = float(width)
        settings["height_mm"] = float(height)
    else:
        raise ValueError("validation_error: create_pdf page.size must be a string preset or a custom object.")

    orientation = str(spec.get("orientation", "portrait")).casefold()
    settings["orientation"] = orientation
    if orientation == "landscape":
        settings["width_mm"], settings["height_mm"] = settings["height_mm"], settings["width_mm"]

    margins_spec = spec.get("margins_mm")
    if isinstance(margins_spec, dict):
        for field_name in ("top", "bottom", "left", "right"):
            value = margins_spec.get(field_name)
            if isinstance(value, (int, float)) and not isinstance(value, bool):
                settings["margins_mm"][field_name] = float(value)
    return settings


def _build_pdf_v2_styles(base_styles: object, style_specs: dict) -> dict[str, object]:
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT, TA_RIGHT
    from reportlab.lib.styles import ParagraphStyle

    align_map = {"left": TA_LEFT, "center": TA_CENTER, "right": TA_RIGHT, "justify": TA_JUSTIFY}

    defaults = {
        "body": {"font": "Helvetica", "size_pt": 11, "line_spacing": 1.15, "space_after_pt": 6},
        "title": {"font": "Helvetica", "size_pt": 22, "bold": True, "alignment": "center", "space_after_pt": 18},
        "heading1": {"font": "Helvetica", "size_pt": 18, "bold": True, "space_before_pt": 18, "space_after_pt": 10},
        "heading2": {"font": "Helvetica", "size_pt": 14, "bold": True, "space_before_pt": 12, "space_after_pt": 8},
        "heading3": {"font": "Helvetica", "size_pt": 12, "bold": True, "space_before_pt": 10, "space_after_pt": 6},
        "quote": {"font": "Helvetica-Oblique", "size_pt": 11, "italic": True, "line_spacing": 1.15, "space_after_pt": 6},
        "blockquote": {"font": "Helvetica-Oblique", "size_pt": 11, "italic": True, "line_spacing": 1.15, "space_after_pt": 6},
    }

    result: dict[str, object] = {}
    heading1_breaks_before = False
    for name, default_spec in defaults.items():
        merged = {**default_spec, **(style_specs.get(name) if isinstance(style_specs.get(name), dict) else {})}
        font_name = _pdf_v2_resolve_font_name(str(merged.get("font", "Helvetica")), bold=bool(merged.get("bold")), italic=bool(merged.get("italic")))
        size_pt = float(merged.get("size_pt", 11))
        line_spacing = float(merged.get("line_spacing", 1.15))
        alignment = align_map.get(str(merged.get("alignment", "left")).casefold(), TA_LEFT)
        result[name] = ParagraphStyle(
            f"PLwC-{name}",
            fontName=font_name,
            fontSize=size_pt,
            leading=size_pt * line_spacing,
            alignment=alignment,
            spaceBefore=float(merged.get("space_before_pt", 0)),
            spaceAfter=float(merged.get("space_after_pt", 0)),
            firstLineIndent=float(merged.get("first_line_indent_mm", 0)) * 2.834645669,  # mm to pt
            underlineProportion=0.04 if merged.get("underline") else 0.0,
            textColor=base_styles["Normal"].textColor,
        )
        if name == "heading1" and merged.get("page_break_before"):
            heading1_breaks_before = True
    result["heading1_breaks_before"] = heading1_breaks_before
    return result


def _pdf_v2_resolve_font_name(font: str, *, bold: bool, italic: bool) -> str:
    """Map (font, bold, italic) to a reportlab built-in font name.

    Only PDF-safe built-in fonts are used. External font files are never
    embedded, bundled or fetched.
    """
    candidate = font.strip() or "Helvetica"
    helvetica = candidate.casefold() in {"helvetica", "arial", "sans-serif"}
    times = candidate.casefold() in {"times", "times-roman", "times new roman", "serif"}
    courier = candidate.casefold() in {"courier", "courier new", "monospace"}
    if helvetica or (not times and not courier):
        if bold and italic:
            return "Helvetica-BoldOblique"
        if bold:
            return "Helvetica-Bold"
        if italic:
            return "Helvetica-Oblique"
        return "Helvetica"
    if times:
        if bold and italic:
            return "Times-BoldItalic"
        if bold:
            return "Times-Bold"
        if italic:
            return "Times-Italic"
        return "Times-Roman"
    # courier
    if bold and italic:
        return "Courier-BoldOblique"
    if bold:
        return "Courier-Bold"
    if italic:
        return "Courier-Oblique"
    return "Courier"


def _apply_pdf_v2_style_overrides(base_style: object, overrides: dict[str, object], align_map: dict) -> object:
    from reportlab.lib.styles import ParagraphStyle

    derived = ParagraphStyle(
        f"{base_style.name}-override",
        parent=base_style,
    )
    if "font" in overrides:
        derived.fontName = _pdf_v2_resolve_font_name(
            str(overrides["font"]),
            bold=bool(overrides.get("bold", False)),
            italic=bool(overrides.get("italic", False)),
        )
    if "size_pt" in overrides:
        derived.fontSize = float(overrides["size_pt"])
        line_spacing = float(overrides.get("line_spacing", 1.15))
        derived.leading = derived.fontSize * line_spacing
    if "line_spacing" in overrides:
        derived.leading = derived.fontSize * float(overrides["line_spacing"])
    if "alignment" in overrides:
        derived.alignment = align_map.get(str(overrides["alignment"]).casefold(), derived.alignment)
    if "space_before_pt" in overrides:
        derived.spaceBefore = float(overrides["space_before_pt"])
    if "space_after_pt" in overrides:
        derived.spaceAfter = float(overrides["space_after_pt"])
    if "first_line_indent_mm" in overrides:
        derived.firstLineIndent = float(overrides["first_line_indent_mm"]) * 2.834645669
    return derived


def _resolve_pdf_v2_image_path(raw_path: str) -> Path:
    text = raw_path.strip()
    if re.match(r"^[A-Za-z][A-Za-z0-9+.\-]+:", text):
        raise ValueError("validation_error: PDF image path must not be a URL.")
    if text.startswith("\\\\") or text.startswith("//"):
        raise ValueError("validation_error: PDF image path must not be a UNC network path.")
    candidate = Path(text)
    if candidate.is_absolute():
        raise ValueError("validation_error: PDF image path must be workspace-relative.")
    parts = [part for part in text.replace("\\", "/").split("/") if part]
    if any(part == ".." for part in parts):
        raise ValueError("validation_error: PDF image path must not contain parent traversal.")
    if any(part.casefold() in PDF_V2_PROTECTED_SEGMENTS for part in parts):
        raise ValueError("validation_error: PDF image path must not target profile or governance directories.")
    if candidate.suffix.casefold() not in PDF_V2_IMAGE_SUFFIXES:
        raise ValueError("validation_error: PDF image path must use .png, .jpg or .jpeg.")
    path = WORK_ROOT / candidate
    resolved = path.resolve(strict=True)
    work_root = WORK_ROOT.resolve(strict=True)
    if resolved != work_root and work_root not in resolved.parents:
        raise ValueError("validation_error: PDF image path must stay under /work.")
    if path.is_symlink() or resolved.is_symlink():
        raise ValueError("validation_error: PDF image symlinks are not supported.")
    if not resolved.is_file():
        raise ValueError("validation_error: PDF image path must be a file.")
    if resolved.stat().st_size > PDF_V2_MAX_IMAGE_BYTES:
        raise ValueError(f"PDF image file exceeds limit of {PDF_V2_MAX_IMAGE_BYTES} bytes.")
    try:
        from PIL import Image as PILImage

        with PILImage.open(resolved) as image:
            if image.format not in {"PNG", "JPEG"}:
                raise ValueError("validation_error: PDF image must be PNG or JPEG.")
            image.verify()
    except ValueError:
        raise
    except Exception as exc:  # pragma: no cover - delegated to PIL
        raise ValueError("validation_error: PDF image is not a readable PNG or JPEG.") from exc
    return resolved


def _validate_pdf_v2_worker_content(content: dict[str, object]) -> dict[str, int]:
    unknown_top = sorted(str(key) for key in content if str(key) not in PDF_V2_TOP_LEVEL_FIELDS)
    if unknown_top:
        raise ValueError(f"validation_error: create_pdf unsupported fields: {', '.join(unknown_top)}.")
    schema_version = content.get("schema_version", PDF_V2_SCHEMA_VERSION)
    if not isinstance(schema_version, str) or schema_version not in PDF_V2_SUPPORTED_SCHEMA_VERSIONS:
        raise ValueError(
            "validation_error: create_pdf schema_version must be one of: "
            f"{', '.join(sorted(PDF_V2_SUPPORTED_SCHEMA_VERSIONS))}."
        )
    elements = content.get("content")
    if not isinstance(elements, list) or not elements:
        raise ValueError("validation_error: create_pdf content.content must be a non-empty list.")
    if len(elements) > PDF_V2_MAX_CONTENT_ELEMENTS:
        raise ValueError(f"create_pdf content element count exceeds {PDF_V2_MAX_CONTENT_ELEMENTS}.")

    stats = {
        "content_element_count": 0,
        "paragraph_count": 0,
        "table_count": 0,
        "table_cell_count": 0,
        "image_count": 0,
        "list_item_count": 0,
    }
    for index, element in enumerate(elements, start=1):
        if not isinstance(element, dict):
            raise ValueError(f"validation_error: create_pdf content element {index} must be an object.")
        element_type = element.get("type")
        if not isinstance(element_type, str) or element_type not in PDF_V2_CONTENT_TYPES:
            raise ValueError(
                f"validation_error: create_pdf content element {index}.type must be one of: "
                f"{', '.join(sorted(PDF_V2_CONTENT_TYPES))}."
            )
        if element_type in PDF_V2_TEXT_TYPES:
            stats["paragraph_count"] += 1
        elif element_type in {"bullet_list", "numbered_list"}:
            items = element.get("items")
            if not isinstance(items, list) or not items:
                raise ValueError(
                    f"validation_error: create_pdf {element_type} element {index}.items must be a non-empty list."
                )
            stats["list_item_count"] += len(items)
            stats["paragraph_count"] += len(items)
        elif element_type == "table":
            stats["table_count"] += 1
            if stats["table_count"] > PDF_V2_MAX_TABLES:
                raise ValueError(f"create_pdf table count exceeds {PDF_V2_MAX_TABLES}.")
            rows = element.get("rows")
            if not isinstance(rows, list) or not rows:
                raise ValueError(f"validation_error: create_pdf table element {index}.rows must be a non-empty list.")
            cell_count = sum(len(row) for row in rows if isinstance(row, list))
            stats["table_cell_count"] += cell_count
            if stats["table_cell_count"] > PDF_V2_MAX_TABLE_CELLS:
                raise ValueError(f"create_pdf table cell count exceeds {PDF_V2_MAX_TABLE_CELLS}.")
        elif element_type == "image":
            stats["image_count"] += 1
            if stats["image_count"] > PDF_V2_MAX_IMAGES:
                raise ValueError(f"create_pdf image count exceeds {PDF_V2_MAX_IMAGES}.")
        stats["content_element_count"] += 1
    if stats["paragraph_count"] > PDF_V2_MAX_PARAGRAPHS:
        raise ValueError(f"create_pdf paragraph count exceeds {PDF_V2_MAX_PARAGRAPHS}.")
    return stats


def _pdf_v2_limits() -> dict[str, int]:
    return {
        "max_json_input_bytes": PDF_V2_MAX_JSON_INPUT_BYTES,
        "max_content_elements": PDF_V2_MAX_CONTENT_ELEMENTS,
        "max_paragraphs": PDF_V2_MAX_PARAGRAPHS,
        "max_tables": PDF_V2_MAX_TABLES,
        "max_table_cells": PDF_V2_MAX_TABLE_CELLS,
        "max_images": PDF_V2_MAX_IMAGES,
        "max_image_bytes": PDF_V2_MAX_IMAGE_BYTES,
    }


def _write_json(payload: dict[str, object]) -> None:
    sys.stdout.write(json.dumps(payload, sort_keys=True))
    sys.stdout.write("\n")


if __name__ == "__main__":
    raise SystemExit(main())
